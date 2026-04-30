from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "presentation" / "Selective_Response_Forum_Congestion_Professional.pptx"

BG = RGBColor(248, 249, 252)
NAVY = RGBColor(24, 36, 61)
BLUE = RGBColor(40, 92, 145)
TEAL = RGBColor(33, 124, 122)
ORANGE = RGBColor(214, 118, 39)
GREEN = RGBColor(52, 122, 84)
RED = RGBColor(183, 61, 61)
MUTED = RGBColor(103, 112, 126)
LINE = RGBColor(216, 222, 230)
WHITE = RGBColor(255, 255, 255)
PALE_BLUE = RGBColor(236, 243, 250)
PALE_TEAL = RGBColor(234, 246, 245)
PALE_ORANGE = RGBColor(255, 243, 231)
PALE_GREEN = RGBColor(236, 247, 240)
PALE_RED = RGBColor(252, 238, 238)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    return shp


def rrect(slide, x, y, w, h, fill, line=None, r=0.12):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = r
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    return shp


def txt(slide, x, y, w, h, text, size=20, color=NAVY, bold=False, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, size=17, color=NAVY):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(3)
    return box


def pic(slide, rel, x, y, w, h):
    slide.shapes.add_picture(str(ROOT / rel), x, y, width=w, height=h)


def base(slide, accent=BLUE):
    rect(slide, 0, 0, prs.slide_width, prs.slide_height, BG)
    rect(slide, Inches(0.4), Inches(0.34), Inches(12.55), Inches(0.12), NAVY)
    rect(slide, Inches(0.4), Inches(0.56), Inches(1.55), Inches(0.10), accent)
    rect(slide, Inches(0.4), Inches(6.95), Inches(12.55), Inches(0.02), LINE)


def header(slide, title, subtitle, chip, chip_color):
    txt(slide, Inches(0.58), Inches(0.82), Inches(8.9), Inches(0.42), title, size=26, bold=True)
    txt(slide, Inches(0.60), Inches(1.24), Inches(9.5), Inches(0.22), subtitle, size=14, color=MUTED)
    rrect(slide, Inches(0.60), Inches(1.65), Inches(1.18), Inches(0.24), chip_color, None, 0.22)
    txt(slide, Inches(0.73), Inches(1.71), Inches(0.92), Inches(0.10), chip, size=11, color=WHITE, bold=True)


def foot(slide, refs, n):
    txt(slide, Inches(0.50), Inches(7.0), Inches(10.0), Inches(0.14), refs, size=9, color=MUTED)
    rrect(slide, Inches(12.48), Inches(6.88), Inches(0.34), Inches(0.26), WHITE, LINE, 0.18)
    txt(slide, Inches(12.58), Inches(6.94), Inches(0.12), Inches(0.10), str(n), size=10, bold=True, align=PP_ALIGN.CENTER)


def title_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    rect(s, Inches(0.58), Inches(0.62), Inches(0.18), Inches(5.7), TEAL)
    txt(s, Inches(0.95), Inches(0.95), Inches(6.2), Inches(1.3), "Selective Response for\nGenAI Under Forum Congestion", size=29, color=WHITE, bold=True)
    txt(s, Inches(0.98), Inches(2.72), Inches(4.8), Inches(0.22), "Game Theory 2026 mini project", size=18, color=RGBColor(221, 228, 238))
    txt(s, Inches(0.98), Inches(3.26), Inches(5.8), Inches(0.22), "Dikshant Gupta (24108)   |   Paras Raina (24170)", size=18, color=WHITE)
    txt(s, Inches(0.98), Inches(3.63), Inches(5.5), Inches(0.22), "Professional deck: intuition + model + results + limitations", size=12, color=RGBColor(221, 228, 238))
    rrect(s, Inches(7.25), Inches(1.02), Inches(5.12), Inches(2.00), WHITE, None, 0.08)
    pic(s, "sim/figs/extra/seed_robustness_ST.png", Inches(7.46), Inches(1.18), Inches(4.70), Inches(1.58))
    rrect(s, Inches(7.25), Inches(3.45), Inches(5.12), Inches(2.00), WHITE, None, 0.08)
    pic(s, "sim/figs/extra/collapse_boundary_c1_xi.png", Inches(7.46), Inches(3.61), Inches(4.70), Inches(1.58))
    txt(s, Inches(7.34), Inches(5.86), Inches(4.7), Inches(0.38), "Question: should GenAI answer now, or preserve the forum that creates future knowledge?", size=16, color=WHITE)


def s2_motivation():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, TEAL)
    header(s, "Motivation", "The key tradeoff is immediate convenience versus long-run knowledge creation.", "Motivation", TEAL)
    rrect(s, Inches(0.65), Inches(2.02), Inches(4.15), Inches(4.20), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.28), Inches(2.8), Inches(0.18), "Intuition for class", size=20, bold=True, color=TEAL)
    bullets(s, Inches(0.88), Inches(2.66), Inches(3.45), Inches(2.85), [
        "If AI answers everything, users stop creating human discussions.",
        "Those missing discussions reduce future signal that AI could have learned from.",
        "But if AI defers too much, the human forum gets overloaded and less useful."
    ], size=18)
    rrect(s, Inches(5.05), Inches(2.02), Inches(7.8), Inches(4.20), PALE_TEAL, LINE, 0.08)
    rrect(s, Inches(5.50), Inches(3.05), Inches(1.25), Inches(0.74), WHITE, TEAL, 0.16)
    rrect(s, Inches(7.32), Inches(3.05), Inches(1.25), Inches(0.74), WHITE, ORANGE, 0.16)
    rrect(s, Inches(9.14), Inches(3.05), Inches(1.35), Inches(0.74), WHITE, GREEN, 0.16)
    rrect(s, Inches(11.06), Inches(3.05), Inches(1.10), Inches(0.74), WHITE, RED, 0.16)
    txt(s, Inches(5.84), Inches(3.30), Inches(0.55), Inches(0.12), "Users", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(7.67), Inches(3.30), Inches(0.55), Inches(0.12), "GenAI", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(9.48), Inches(3.30), Inches(0.62), Inches(0.12), "Forum", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(11.34), Inches(3.30), Inches(0.40), Inches(0.12), "C_t", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(6.83), Inches(3.21), Inches(0.22), Inches(0.10), "->", size=22, color=TEAL, bold=True)
    txt(s, Inches(8.63), Inches(3.21), Inches(0.22), Inches(0.10), "->", size=22, color=ORANGE, bold=True)
    txt(s, Inches(10.55), Inches(3.21), Inches(0.22), Inches(0.10), "->", size=22, color=GREEN, bold=True)
    txt(s, Inches(5.55), Inches(4.55), Inches(6.9), Inches(0.60), "So selective response is not only a service decision. It is a system-design decision.", size=21, bold=True, align=PP_ALIGN.CENTER)
    foot(s, "Refs: Taitler et al. (2025); Hardin (1968); Ostrom (1990)", 2)


def s3_examples():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, ORANGE)
    header(s, "Why This Problem Matters", "The same mechanism appears in education, developer communities, and product support.", "Examples", ORANGE)
    cards = [
        ("Course forums", "Too much AI can reduce the student explanations that help future batches.", PALE_ORANGE),
        ("Developer forums", "Overloaded expert communities create slower and lower-quality answers.", PALE_BLUE),
        ("Product communities", "Short-run automation may weaken the long-run human knowledge base.", PALE_GREEN),
    ]
    for i, (title, body, fill) in enumerate(cards):
        x = 0.7 + 4*i
        rrect(s, Inches(x), Inches(2.12), Inches(3.55), Inches(4.15), fill, LINE, 0.08)
        txt(s, Inches(x + 0.22), Inches(2.40), Inches(2.8), Inches(0.18), title, size=19, bold=True, color=ORANGE if i == 0 else BLUE if i == 1 else GREEN)
        txt(s, Inches(x + 0.22), Inches(3.02), Inches(2.95), Inches(1.5), body, size=18)
        txt(s, Inches(x + 0.22), Inches(5.30), Inches(2.9), Inches(0.35), "Human discussion here is a productive resource.", size=14, color=MUTED)
    foot(s, "Refs: Horvitz (1999); Amershi et al. (2014); Bansal et al. (2021)", 3)


def s4_literature():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, BLUE)
    header(s, "Literature Review", "Our project combines three strands of prior work.", "Literature", BLUE)
    entries = [
        ("Selective response", PALE_BLUE, BLUE, "Chow (1970)\nEl-Yaniv (2010)", "When should the model abstain?"),
        ("Learning to defer", PALE_ORANGE, ORANGE, "Madras et al. (2018)\nMozannar and Sontag (2020)", "When should a human handle the task?"),
        ("Congestion / commons", PALE_GREEN, GREEN, "Hardin (1968)\nOstrom (1990)\nKleinrock (1975)", "What happens when shared systems are overloaded?"),
    ]
    for i, (title, fill, accent, refs, question) in enumerate(entries):
        x = 0.7 + 4*i
        rrect(s, Inches(x), Inches(2.10), Inches(3.55), Inches(4.15), fill, LINE, 0.08)
        txt(s, Inches(x + 0.22), Inches(2.36), Inches(2.9), Inches(0.18), title, size=19, bold=True, color=accent)
        txt(s, Inches(x + 0.22), Inches(2.95), Inches(2.8), Inches(0.75), refs, size=17)
        txt(s, Inches(x + 0.22), Inches(4.45), Inches(3.0), Inches(0.85), question, size=18, color=MUTED)
    rrect(s, Inches(1.15), Inches(6.27), Inches(10.95), Inches(0.40), WHITE, LINE, 0.08)
    txt(s, Inches(1.42), Inches(6.38), Inches(10.4), Inches(0.12), "Gap we address: prior selective-response work does not model a congestible human forum with endogenous capacity.", size=17, bold=True, align=PP_ALIGN.CENTER)
    foot(s, "Refs: Chow (1970); El-Yaniv (2010); Madras et al. (2018); Mozannar and Sontag (2020); Hardin (1968); Ostrom (1990)", 4)


def s5_basepaper():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, ORANGE)
    header(s, "Base Paper And Its Limitation", "This part is prior work; our project starts from its missing systems detail.", "Prior work", ORANGE)
    rrect(s, Inches(0.65), Inches(2.08), Inches(4.05), Inches(4.18), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.34), Inches(2.6), Inches(0.18), "Base paper message", size=20, bold=True, color=ORANGE)
    bullets(s, Inches(0.88), Inches(2.72), Inches(3.35), Inches(2.55), [
        "Selective response can improve long-run learning.",
        "Reason: some unanswered queries create future human knowledge.",
        "So always-answering is not automatically optimal."
    ], size=18)
    rrect(s, Inches(5.0), Inches(2.08), Inches(2.9), Inches(4.18), PALE_ORANGE, LINE, 0.08)
    txt(s, Inches(5.28), Inches(2.34), Inches(1.8), Inches(0.18), "Simple intuition", size=20, bold=True, color=ORANGE)
    rrect(s, Inches(5.58), Inches(3.08), Inches(1.75), Inches(0.72), WHITE, ORANGE, 0.18)
    rrect(s, Inches(5.58), Inches(4.58), Inches(1.75), Inches(0.72), WHITE, GREEN, 0.18)
    txt(s, Inches(5.95), Inches(3.33), Inches(1.0), Inches(0.12), "Always answer", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(5.86), Inches(4.83), Inches(1.15), Inches(0.12), "Selective response", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(5.16), Inches(5.62), Inches(2.55), Inches(0.35), "Less answering can sometimes create more future learning.", size=16, bold=True, align=PP_ALIGN.CENTER)
    rrect(s, Inches(8.2), Inches(2.08), Inches(4.65), Inches(4.18), PALE_RED, LINE, 0.08)
    txt(s, Inches(8.48), Inches(2.34), Inches(2.9), Inches(0.18), "What the base paper leaves out", size=20, bold=True, color=RED)
    bullets(s, Inches(8.44), Inches(2.72), Inches(3.95), Inches(2.60), [
        "The forum is not modeled as capacity-limited.",
        "There is no overload damage or collapse dynamic.",
        "It does not isolate volume effects from composition effects."
    ], size=18)
    foot(s, "Refs: Taitler et al. (2025)", 5)


def s6_model():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, TEAL)
    header(s, "Our Model: Intuition And Mathematics", "We keep the simulator reduced-form, but we make the mechanism explicit.", "Model", TEAL)
    rrect(s, Inches(0.65), Inches(2.06), Inches(4.10), Inches(4.20), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.32), Inches(2.4), Inches(0.18), "Economic intuition", size=20, bold=True, color=TEAL)
    bullets(s, Inches(0.88), Inches(2.68), Inches(3.4), Inches(2.80), [
        "Users choose between GenAI and the forum.",
        "Forum output creates future knowledge stock S_t.",
        "If load exceeds capacity C_t, useful throughput is discounted.",
        "Repeated overload reduces future capacity."
    ], size=18)
    rrect(s, Inches(4.98), Inches(2.06), Inches(7.87), Inches(4.20), PALE_TEAL, LINE, 0.08)
    txt(s, Inches(5.25), Inches(2.28), Inches(2.2), Inches(0.18), "Core equations", size=20, bold=True, color=TEAL)
    txt(s, Inches(5.25), Inches(2.78), Inches(7.1), Inches(0.46), "m_t = 1 - p_t^answered", size=20, color=NAVY, font="Consolas")
    txt(s, Inches(5.25), Inches(3.22), Inches(7.1), Inches(0.46), "G_t = kappa * throughput_t * novelty_factor_t", size=20, color=NAVY, font="Consolas")
    txt(s, Inches(5.25), Inches(3.66), Inches(7.1), Inches(0.46), "S_{t+1} = S_t + G_t", size=20, color=NAVY, font="Consolas")
    txt(s, Inches(5.25), Inches(4.10), Inches(7.1), Inches(0.46), "C_{t+1} = max{0, (1-delta)C_t + eta R - xi max(m_t - C_t, 0)}", size=18, color=NAVY, font="Consolas")
    txt(s, Inches(5.25), Inches(5.02), Inches(6.9), Inches(0.62), "Professor-facing intuition: xi controls forum fragility, lambda controls how useful overloaded traffic remains, and rho controls how much novelty matters.", size=16, color=MUTED)
    foot(s, "Refs: model summarized from paper/main.tex and simulator", 6)


def s7_design():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, BLUE)
    header(s, "Experiment Design", "We test both the base intuition and our new congestion-aware mechanisms.", "Design", BLUE)
    rrect(s, Inches(0.65), Inches(2.08), Inches(4.0), Inches(4.18), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.34), Inches(2.0), Inches(0.18), "Policies", size=20, bold=True, color=BLUE)
    bullets(s, Inches(0.88), Inches(2.72), Inches(3.3), Inches(2.65), [
        "Uniform throttling",
        "Targeted novelty policy",
        "Capacity-aware baseline",
        "Novelty + capacity baseline"
    ], size=18)
    rrect(s, Inches(4.9), Inches(2.08), Inches(3.7), Inches(4.18), PALE_BLUE, LINE, 0.08)
    txt(s, Inches(5.18), Inches(2.34), Inches(2.1), Inches(0.18), "Regimes", size=20, bold=True, color=BLUE)
    bullets(s, Inches(5.14), Inches(2.72), Inches(3.0), Inches(2.20), [
        "Stable: C1 = 0.85, xi = 0.04",
        "Mid: C1 = 0.60, xi = 0.10",
        "Collapse: C1 = 0.35, xi = 0.35",
        "30 novelty seeds"
    ], size=18)
    rrect(s, Inches(8.85), Inches(2.08), Inches(4.0), Inches(4.18), PALE_GREEN, LINE, 0.08)
    txt(s, Inches(9.12), Inches(2.34), Inches(2.2), Inches(0.18), "Metrics", size=20, bold=True, color=GREEN)
    bullets(s, Inches(9.08), Inches(2.72), Inches(3.25), Inches(2.35), [
        "Final learning S_T",
        "Minimum capacity",
        "Overload time",
        "Average answered share",
        "Welfare and revenue proxies"
    ], size=18)
    foot(s, "Refs: sim/results/summary.md; sim/results/extra/summary_extra.md", 7)


def s8_results():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, ORANGE)
    header(s, "Main Results", "Novelty-aware routing improves learning and often reduces overload.", "Results", ORANGE)
    rrect(s, Inches(0.65), Inches(2.08), Inches(3.4), Inches(4.18), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.34), Inches(2.2), Inches(0.18), "Headline numbers", size=20, bold=True, color=ORANGE)
    bullets(s, Inches(0.88), Inches(2.72), Inches(2.8), Inches(2.65), [
        "Stable: 26.15 -> 36.29",
        "Mid: 19.60 -> 31.56",
        "Collapse: 6.68 -> 9.54",
        "Mid overload: 40 -> 25"
    ], size=18)
    txt(s, Inches(0.92), Inches(5.40), Inches(2.4), Inches(0.36), "Interpretation: the forum receives more useful traffic.", size=15, color=MUTED)
    rrect(s, Inches(4.32), Inches(2.08), Inches(4.2), Inches(4.18), PALE_BLUE, LINE, 0.08)
    pic(s, "sim/figs/extra/seed_robustness_ST.png", Inches(4.50), Inches(2.26), Inches(3.85), Inches(3.25))
    rrect(s, Inches(8.78), Inches(2.08), Inches(4.07), Inches(4.18), PALE_BLUE, LINE, 0.08)
    pic(s, "sim/figs/extra/seed_robustness_overload.png", Inches(8.96), Inches(2.26), Inches(3.72), Inches(3.25))
    foot(s, "Refs: seed_robustness.csv", 8)


def s9_validation():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, RED)
    header(s, "Mechanism Validation", "We explicitly test whether the gain is composition-driven.", "Validation", RED)
    rrect(s, Inches(0.65), Inches(2.08), Inches(3.2), Inches(4.18), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.34), Inches(2.1), Inches(0.18), "Two key checks", size=20, bold=True, color=RED)
    bullets(s, Inches(0.88), Inches(2.72), Inches(2.6), Inches(2.55), [
        "Matched answered-share control",
        "rho = 0 ablation",
        "If both pass, the novelty mechanism is credible"
    ], size=18)
    rrect(s, Inches(4.05), Inches(2.08), Inches(4.05), Inches(4.18), PALE_RED, LINE, 0.08)
    pic(s, "sim/figs/extra/matched_answered_compare_ST.png", Inches(4.23), Inches(2.26), Inches(3.70), Inches(3.25))
    rrect(s, Inches(8.40), Inches(2.08), Inches(4.45), Inches(4.18), PALE_RED, LINE, 0.08)
    pic(s, "sim/figs/extra/rho_ablation_gap_ST.png", Inches(8.58), Inches(2.26), Inches(4.10), Inches(3.25))
    foot(s, "Refs: matched_answered_compare.csv; rho_ablation.csv", 9)


def s10_tipping():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, GREEN)
    header(s, "Tipping And Stronger Baselines", "Congestion feedback helps, but novelty + capacity works best in the mid regime.", "Phase view", GREEN)
    rrect(s, Inches(0.65), Inches(2.08), Inches(3.2), Inches(4.18), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.34), Inches(2.1), Inches(0.18), "Policy ranking", size=20, bold=True, color=GREEN)
    bullets(s, Inches(0.88), Inches(2.72), Inches(2.55), Inches(2.55), [
        "Uniform: 19.60",
        "Capacity-aware: 22.52",
        "Targeted novelty: 31.56",
        "Novelty + capacity: 33.23"
    ], size=18)
    rrect(s, Inches(4.05), Inches(2.08), Inches(4.15), Inches(4.18), PALE_GREEN, LINE, 0.08)
    pic(s, "sim/figs/extra/collapse_boundary_c1_xi.png", Inches(4.23), Inches(2.26), Inches(3.80), Inches(3.25))
    rrect(s, Inches(8.50), Inches(2.08), Inches(4.35), Inches(4.18), PALE_GREEN, LINE, 0.08)
    pic(s, "sim/figs/extra/policy_compare_extended_ST.png", Inches(8.68), Inches(2.26), Inches(4.00), Inches(3.25))
    foot(s, "Refs: collapse_boundary_c1_xi.csv; policy_compare_extended.csv", 10)


def s11_summary():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, ORANGE)
    header(s, "Summary And Future Work", "End by separating clearly what comes from the base paper and what comes from us.", "Summary", ORANGE)
    rrect(s, Inches(0.65), Inches(2.08), Inches(4.2), Inches(4.18), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.34), Inches(2.2), Inches(0.18), "Summary", size=20, bold=True, color=ORANGE)
    bullets(s, Inches(0.88), Inches(2.72), Inches(3.5), Inches(2.55), [
        "Base paper: selective response can improve long-run learning.",
        "Our contribution: forums are congestible and capacity-limited.",
        "Main lesson: defer the right traffic, not just more traffic."
    ], size=18)
    rrect(s, Inches(5.10), Inches(2.08), Inches(3.55), Inches(4.18), PALE_ORANGE, LINE, 0.08)
    txt(s, Inches(5.38), Inches(2.34), Inches(1.9), Inches(0.18), "Future work", size=20, bold=True, color=ORANGE)
    bullets(s, Inches(5.34), Inches(2.72), Inches(2.85), Inches(2.35), [
        "Use real forum data for calibration.",
        "Learn routing policies.",
        "Model multiple forums and heterogeneous experts."
    ], size=18)
    rrect(s, Inches(8.95), Inches(2.08), Inches(3.90), Inches(4.18), PALE_BLUE, LINE, 0.08)
    txt(s, Inches(9.22), Inches(2.34), Inches(1.8), Inches(0.18), "References", size=20, bold=True, color=BLUE)
    txt(s, Inches(9.18), Inches(2.72), Inches(3.15), Inches(2.8), "Taitler et al. (2025)\nChow (1970)\nEl-Yaniv (2010)\nMadras et al. (2018)\nMozannar and Sontag (2020)\nHardin (1968)\nOstrom (1990)\nVickrey (1969)", size=16)
    foot(s, "Refs: full bibliography in refs.bib", 11)


def thanks():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    rrect(s, Inches(0.95), Inches(1.05), Inches(5.15), Inches(4.75), WHITE, None, 0.08)
    txt(s, Inches(1.30), Inches(1.65), Inches(3.2), Inches(0.35), "Thank You", size=34, bold=True)
    txt(s, Inches(1.32), Inches(2.45), Inches(2.0), Inches(0.22), "Questions?", size=24, color=TEAL, bold=True)
    txt(s, Inches(1.32), Inches(3.15), Inches(4.2), Inches(0.22), "Dikshant Gupta (24108)", size=18)
    txt(s, Inches(1.32), Inches(3.50), Inches(4.2), Inches(0.22), "Paras Raina (24170)", size=18)
    rrect(s, Inches(7.05), Inches(1.28), Inches(5.15), Inches(3.75), WHITE, None, 0.08)
    pic(s, "sim/figs/extra/policy_compare_extended_ST.png", Inches(7.28), Inches(1.52), Inches(4.7), Inches(2.95))
    txt(s, Inches(7.30), Inches(5.28), Inches(4.3), Inches(0.25), "Final takeaway: novelty-aware routing is strongest when combined with congestion response.", size=16, color=WHITE)


OUT.parent.mkdir(exist_ok=True)
title_slide()
s2_motivation()
s3_examples()
s4_literature()
s5_basepaper()
s6_model()
s7_design()
s8_results()
s9_validation()
s10_tipping()
s11_summary()
thanks()
prs.save(OUT)
print(OUT)
