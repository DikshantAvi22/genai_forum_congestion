from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "presentation"
PPTX_OUT = OUT_DIR / "Selective_Response_Forum_Congestion_Presentation_Editable.pptx"

BLUE = RGBColor(31, 78, 121)
MID_BLUE = RGBColor(68, 114, 196)
ACCENT = RGBColor(217, 119, 6)
GREEN = RGBColor(22, 101, 52)
RED = RGBColor(185, 28, 28)
TEXT = RGBColor(20, 33, 61)
MUTED = RGBColor(95, 107, 122)
BG = RGBColor(246, 248, 252)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(234, 241, 251)
LIGHT_ORANGE = RGBColor(255, 243, 232)
LIGHT_GREEN = RGBColor(237, 248, 241)
LIGHT_RED = RGBColor(253, 236, 236)
GRID = RGBColor(217, 227, 240)


def add_full_bg(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.35), Inches(12.9), Inches(0.10))
    top.fill.solid()
    top.fill.fore_color.rgb = BLUE
    top.line.fill.background()
    bot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(7.02), Inches(12.9), Inches(0.02))
    bot.fill.solid()
    bot.fill.fore_color.rgb = GRID
    bot.line.fill.background()


def add_panel(slide, x, y, w, h, fill=WHITE, line=GRID, radius=0.12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    return shape


def add_text(slide, x, y, w, h, text, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT, font_name="Aptos", margin=0.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_kicker(slide, x, y, text, fill, width=1.15):
    add_badge(slide, x, y, Inches(width), Inches(0.24), text, fill)


def add_bullets(slide, x, y, w, h, items, size=18, color=TEXT, bullet_color=MID_BLUE):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.bullet = True
        p.space_after = Pt(6)
    return box


def add_badge(slide, x, y, w, h, text, fill, text_color=WHITE):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.25
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    add_text(slide, x + Inches(0.08), y + Inches(0.01), w - Inches(0.16), h - Inches(0.02), text, size=13, color=text_color, bold=True)


def add_footer(slide, refs, n):
    add_text(slide, Inches(0.50), Inches(6.98), Inches(9.2), Inches(0.25), refs, size=10, color=MUTED)
    add_badge(slide, Inches(12.55), Inches(6.92), Inches(0.32), Inches(0.28), str(n), LIGHT_BLUE, BLUE)


def add_picture(slide, rel_path, x, y, w, h):
    slide.shapes.add_picture(str(ROOT / rel_path), x, y, width=w, height=h)


def slide_title(slide, title, subtitle):
    add_text(slide, Inches(0.55), Inches(0.70), Inches(7.4), Inches(0.85), title, size=30, color=TEXT, bold=True)
    add_text(slide, Inches(0.55), Inches(1.52), Inches(8.5), Inches(0.35), subtitle, size=16, color=MUTED)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_text(slide, Inches(0.55), Inches(0.60), Inches(5.6), Inches(1.3), "Selective Response for\nGenAI Under Forum\nCongestion", size=28, color=TEXT, bold=True)
    add_text(slide, Inches(0.55), Inches(1.88), Inches(3.6), Inches(0.3), "Game Theory 2026 mini project", size=16, color=MUTED)
    add_badge(slide, Inches(0.55), Inches(2.28), Inches(2.45), Inches(0.24), "Base paper: Taitler et al. (2025)", BLUE)
    add_badge(slide, Inches(3.03), Inches(2.28), Inches(3.28), Inches(0.24), "Our work: congestion + mechanism tests", ACCENT)
    add_text(slide, Inches(0.55), Inches(2.70), Inches(5.5), Inches(0.3), "Dikshant Gupta (24108)   |   Paras Raina (24170)", size=18, color=TEXT)
    add_text(slide, Inches(0.55), Inches(2.97), Inches(4.0), Inches(0.20), "Speaker split: Dikshant slides 1-4, Paras slides 5-9", size=12, color=MUTED)

    add_panel(slide, Inches(0.55), Inches(3.18), Inches(5.75), Inches(3.02))
    add_text(slide, Inches(0.78), Inches(3.38), Inches(2.2), Inches(0.3), "Talk in one line", size=20, color=BLUE, bold=True)
    add_bullets(
        slide,
        Inches(0.75),
        Inches(3.72),
        Inches(5.1),
        Inches(2.18),
        [
            "Selective response can help long-run learning, but forums are not infinite-capacity resources.",
            "We replicate the base paper's selective-response intuition and extend it with forum congestion and endogenous capacity.",
            "Main finding: novelty-aware routing outperforms uniform throttling, and the gain survives stricter controls.",
        ],
        size=18,
    )

    add_panel(slide, Inches(6.55), Inches(1.70), Inches(3.05), Inches(2.08))
    add_picture(slide, "sim/figs/extra/seed_robustness_ST.png", Inches(6.75), Inches(1.82), Inches(2.65), Inches(1.45))
    add_text(slide, Inches(6.68), Inches(3.48), Inches(2.6), Inches(0.22), "Learning gains across regimes", size=11, color=MUTED)
    add_panel(slide, Inches(9.85), Inches(1.70), Inches(3.05), Inches(2.08))
    add_picture(slide, "sim/figs/extra/collapse_boundary_c1_xi.png", Inches(10.05), Inches(1.82), Inches(2.65), Inches(1.45))
    add_text(slide, Inches(9.98), Inches(3.48), Inches(2.7), Inches(0.30), "Congestion creates a collapse boundary", size=11, color=MUTED)
    add_panel(slide, Inches(6.55), Inches(4.08), Inches(6.35), Inches(2.12))
    add_picture(slide, "sim/figs/timeseries_collapse.png", Inches(6.75), Inches(4.18), Inches(5.95), Inches(1.45))
    add_text(slide, Inches(6.68), Inches(5.95), Inches(2.6), Inches(0.22), "Collapse regime example", size=11, color=MUTED)
    add_footer(slide, "Refs: Taitler et al. (2025); Ostrom (1990); Vickrey (1969)", 1)


def problem_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    slide_title(slide, "Problem And Why It Matters", "Routing users between GenAI and human forums changes both current service and future knowledge creation.")
    add_kicker(slide, Inches(0.55), Inches(1.95), "Motivation", BLUE)
    add_panel(slide, Inches(0.55), Inches(2.00), Inches(5.5), Inches(4.20))
    add_text(slide, Inches(0.78), Inches(2.20), Inches(2.4), Inches(0.25), "Core problem", size=20, color=BLUE, bold=True)
    add_bullets(
        slide,
        Inches(0.74),
        Inches(2.50),
        Inches(4.95),
        Inches(1.80),
        [
            "If GenAI answers everything, fewer high-value forum discussions happen, so future training signal shrinks.",
            "If GenAI defers too much, forum load rises, response quality drops, and expert capacity can collapse.",
            "The planner must manage both traffic volume and traffic composition.",
        ],
        size=18,
    )
    add_panel(slide, Inches(0.82), Inches(4.35), Inches(4.95), Inches(0.48), fill=LIGHT_BLUE)
    add_text(slide, Inches(1.00), Inches(4.49), Inches(1.4), Inches(0.18), "Course forums", size=13, bold=True)
    add_text(slide, Inches(2.50), Inches(4.48), Inches(3.0), Inches(0.20), "Hard questions create reusable explanations for later students.", size=12, color=MUTED)
    add_panel(slide, Inches(0.82), Inches(4.92), Inches(4.95), Inches(0.48))
    add_text(slide, Inches(1.00), Inches(5.06), Inches(1.4), Inches(0.18), "Developer communities", size=13, bold=True)
    add_text(slide, Inches(2.50), Inches(5.05), Inches(3.0), Inches(0.20), "Stack Overflow style forums can be swamped by redirected users.", size=12, color=MUTED)
    add_panel(slide, Inches(0.82), Inches(5.49), Inches(4.95), Inches(0.48))
    add_text(slide, Inches(1.00), Inches(5.63), Inches(1.4), Inches(0.18), "Product support", size=13, bold=True)
    add_text(slide, Inches(2.50), Inches(5.62), Inches(3.0), Inches(0.20), "Short-term automation can reduce the human knowledge pipeline.", size=12, color=MUTED)

    add_panel(slide, Inches(6.25), Inches(2.00), Inches(6.65), Inches(4.20))
    add_badge(slide, Inches(6.45), Inches(2.25), Inches(1.0), Inches(0.48), "Users", LIGHT_BLUE, TEXT)
    add_badge(slide, Inches(8.05), Inches(2.25), Inches(1.0), Inches(0.48), "GenAI", LIGHT_ORANGE, TEXT)
    add_badge(slide, Inches(9.70), Inches(2.25), Inches(1.1), Inches(0.48), "Forum", LIGHT_GREEN, TEXT)
    add_badge(slide, Inches(11.40), Inches(2.25), Inches(0.9), Inches(0.48), "C_t", LIGHT_RED, TEXT)
    add_text(slide, Inches(6.40), Inches(2.80), Inches(1.1), Inches(0.20), "choice share", size=12, color=MUTED)
    add_text(slide, Inches(8.00), Inches(2.80), Inches(1.2), Inches(0.20), "answered share", size=12, color=MUTED)
    add_text(slide, Inches(9.55), Inches(2.80), Inches(1.5), Inches(0.20), "throughput + novelty", size=12, color=MUTED)
    add_text(slide, Inches(11.35), Inches(2.80), Inches(0.9), Inches(0.20), "capacity", size=12, color=MUTED)
    add_text(slide, Inches(7.45), Inches(2.36), Inches(0.5), Inches(0.25), "->", size=28, color=MID_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(9.10), Inches(2.36), Inches(0.5), Inches(0.25), "->", size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(10.95), Inches(2.36), Inches(0.5), Inches(0.25), "->", size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(10.65), Inches(3.35), Inches(2.0), Inches(0.25), "Overload damages future capacity", size=14, color=RED, bold=True)
    add_panel(slide, Inches(7.10), Inches(4.00), Inches(4.75), Inches(1.10), fill=LIGHT_RED, line=RED)
    add_bullets(
        slide,
        Inches(7.28),
        Inches(4.20),
        Inches(4.35),
        Inches(0.70),
        [
            "Load above capacity reduces useful throughput.",
            "Repeated overload damages future capacity.",
            "This creates tipping behavior instead of smooth decline.",
        ],
        size=13,
        bullet_color=RED,
    )
    add_picture(slide, "sim/figs/timeseries_stable.png", Inches(6.55), Inches(5.35), Inches(2.85), Inches(0.72))
    add_picture(slide, "sim/figs/timeseries_collapse.png", Inches(9.75), Inches(5.35), Inches(2.85), Inches(0.72))
    add_footer(slide, "Refs: Taitler et al. (2025); Hardin (1968); Ostrom (1990); Vickrey (1969)", 2)


def literature_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    slide_title(slide, "Literature Review", "The project combines selective response, human-AI delegation, and congestion economics.")
    add_kicker(slide, Inches(0.55), Inches(1.95), "Context", BLUE, width=0.95)
    add_panel(slide, Inches(0.55), Inches(2.05), Inches(3.6), Inches(3.60), fill=LIGHT_BLUE, line=MID_BLUE)
    add_text(slide, Inches(0.78), Inches(2.25), Inches(3.0), Inches(0.35), "Selective response / reject option", size=19, color=BLUE, bold=True)
    add_text(slide, Inches(0.78), Inches(2.70), Inches(3.0), Inches(2.4), "Chow (1970), El-Yaniv (2010), Geifman and El-Yaniv (2017): when not to answer is part of the decision rule.", size=16)
    add_panel(slide, Inches(4.85), Inches(2.05), Inches(3.6), Inches(3.60), fill=LIGHT_ORANGE, line=ACCENT)
    add_text(slide, Inches(5.08), Inches(2.25), Inches(3.0), Inches(0.35), "Learning to defer / human-AI systems", size=19, color=ACCENT, bold=True)
    add_text(slide, Inches(5.08), Inches(2.70), Inches(3.0), Inches(2.4), "Madras et al. (2018), Mozannar and Sontag (2020), Kleinberg et al. (2018), Amershi et al. (2014): defer to humans when they add value.", size=16)
    add_panel(slide, Inches(9.15), Inches(2.05), Inches(3.75), Inches(3.60), fill=LIGHT_GREEN, line=GREEN)
    add_text(slide, Inches(9.38), Inches(2.25), Inches(3.1), Inches(0.35), "Congestion / shared-resource view", size=19, color=GREEN, bold=True)
    add_text(slide, Inches(9.38), Inches(2.70), Inches(3.1), Inches(2.4), "Hardin (1968), Ostrom (1990), Kleinrock (1975), Vickrey (1969), Roughgarden and Tardos (2006): overloaded shared systems degrade.", size=16)
    add_panel(slide, Inches(2.15), Inches(5.90), Inches(8.85), Inches(0.60))
    add_text(slide, Inches(2.35), Inches(6.10), Inches(8.4), Inches(0.22), "Gap addressed here: prior selective-response work models knowledge creation, but not forum congestion and endogenous capacity loss.", size=16, bold=True)
    add_footer(slide, "Refs: Chow (1970); El-Yaniv (2010); Madras et al. (2018); Mozannar and Sontag (2020); Hardin (1968); Ostrom (1990); Kleinrock (1975)", 3)


def base_paper_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    slide_title(slide, "Base Paper: Idea And Limits", "Selective-response intuition comes from the prior paper; congestion modeling does not.")
    add_badge(slide, Inches(0.55), Inches(1.95), Inches(2.25), Inches(0.24), "Replication / prior work", BLUE)
    add_panel(slide, Inches(0.55), Inches(2.05), Inches(4.15), Inches(4.35))
    add_text(slide, Inches(0.80), Inches(2.25), Inches(3.4), Inches(0.30), "What the base paper studies", size=20, color=BLUE, bold=True)
    add_bullets(
        slide,
        Inches(0.78),
        Inches(2.55),
        Inches(3.6),
        Inches(3.4),
        [
            "Question: should GenAI answer every query, or selectively defer some users to human communities?",
            "Key claim: answering less can improve long-run learning because deferred interactions create future human-generated knowledge.",
            "Interpretation: abstention is not only about current accuracy; it also shapes the future data ecosystem.",
        ],
        size=18,
        bullet_color=MID_BLUE,
    )
    add_panel(slide, Inches(5.00), Inches(2.05), Inches(3.65), Inches(4.35), fill=LIGHT_BLUE, line=MID_BLUE)
    add_text(slide, Inches(5.23), Inches(2.25), Inches(2.8), Inches(0.30), "Base paper takeaway", size=20, color=BLUE, bold=True)
    add_panel(slide, Inches(5.55), Inches(3.00), Inches(2.55), Inches(0.78), fill=WHITE, line=MID_BLUE)
    add_text(slide, Inches(5.78), Inches(3.20), Inches(2.0), Inches(0.20), "Always answer", size=16, bold=True)
    add_text(slide, Inches(5.78), Inches(3.47), Inches(2.0), Inches(0.18), "Less human discussion", size=14, color=MUTED)
    add_panel(slide, Inches(5.55), Inches(4.62), Inches(2.55), Inches(0.78), fill=WHITE, line=GREEN)
    add_text(slide, Inches(5.74), Inches(4.82), Inches(2.1), Inches(0.20), "Selective response", size=16, bold=True)
    add_text(slide, Inches(5.74), Inches(5.09), Inches(2.1), Inches(0.18), "More future signal", size=14, color=MUTED)
    add_text(slide, Inches(6.55), Inches(3.88), Inches(0.22), Inches(0.40), "^", size=30, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.23), Inches(5.72), Inches(3.0), Inches(0.40), "Qualitative result: some deferral can beat always-answering in long-run learning.", size=14)
    add_panel(slide, Inches(8.95), Inches(2.05), Inches(3.95), Inches(4.35), fill=LIGHT_RED, line=RED)
    add_text(slide, Inches(9.18), Inches(2.25), Inches(3.3), Inches(0.55), "Why that is not enough for\nour setting", size=20, color=RED, bold=True)
    add_bullets(
        slide,
        Inches(9.15),
        Inches(2.85),
        Inches(3.3),
        Inches(3.2),
        [
            "Forum is treated as an implicit knowledge source, not a congestible shared resource.",
            "No endogenous forum capacity or collapse dynamics.",
            "No test of whether targeted gains survive when answered share is matched.",
            "No mechanism ablation to verify that novelty is the real source of the gain.",
        ],
        size=17,
        bullet_color=RED,
    )
    add_footer(slide, "Refs: Taitler et al. (2025); Chow (1970); El-Yaniv (2010)", 4)


def method_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    slide_title(slide, "Our Original Work: Methodology", "We extend the base idea with congestion dynamics, new baselines, and mechanism-validation experiments.")
    add_badge(slide, Inches(0.55), Inches(1.95), Inches(1.25), Inches(0.24), "Our extension", ACCENT)
    add_panel(slide, Inches(0.55), Inches(2.05), Inches(8.30), Inches(4.35))
    add_text(slide, Inches(0.78), Inches(2.25), Inches(3.0), Inches(0.30), "Extended simulator", size=20, color=ACCENT, bold=True)
    add_panel(slide, Inches(0.88), Inches(2.72), Inches(1.80), Inches(0.78), fill=LIGHT_BLUE, line=MID_BLUE)
    add_panel(slide, Inches(2.85), Inches(2.72), Inches(1.90), Inches(0.78), fill=LIGHT_ORANGE, line=ACCENT)
    add_panel(slide, Inches(5.30), Inches(2.72), Inches(2.00), Inches(0.78), fill=LIGHT_GREEN, line=GREEN)
    add_panel(slide, Inches(7.30), Inches(2.72), Inches(1.25), Inches(0.78), fill=LIGHT_RED, line=RED)
    add_text(slide, Inches(1.28), Inches(3.00), Inches(1.0), Inches(0.18), "Users", size=17, bold=True)
    add_text(slide, Inches(3.35), Inches(3.00), Inches(1.0), Inches(0.18), "GenAI", size=17, bold=True)
    add_text(slide, Inches(5.85), Inches(3.00), Inches(1.0), Inches(0.18), "Forum", size=17, bold=True)
    add_text(slide, Inches(7.76), Inches(3.00), Inches(0.5), Inches(0.18), "C_t", size=17, bold=True)
    add_text(slide, Inches(1.14), Inches(3.28), Inches(1.3), Inches(0.16), "choice share", size=12, color=MUTED)
    add_text(slide, Inches(3.18), Inches(3.28), Inches(1.3), Inches(0.16), "answered share", size=12, color=MUTED)
    add_text(slide, Inches(5.53), Inches(3.28), Inches(1.6), Inches(0.16), "throughput + novelty", size=12, color=MUTED)
    add_text(slide, Inches(7.58), Inches(3.28), Inches(0.8), Inches(0.16), "capacity", size=12, color=MUTED)
    add_text(slide, Inches(2.55), Inches(2.98), Inches(0.28), Inches(0.20), "->", size=24, color=MID_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4.78), Inches(2.98), Inches(0.28), Inches(0.20), "->", size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(8.32), Inches(2.98), Inches(0.28), Inches(0.20), "->", size=24, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.70), Inches(3.75), Inches(2.0), Inches(0.20), "Overload damages future capacity", size=13, color=RED, bold=True)
    add_panel(slide, Inches(0.88), Inches(4.55), Inches(6.95), Inches(1.35))
    add_text(slide, Inches(1.05), Inches(4.72), Inches(2.5), Inches(0.25), "Key policies compared", size=20, bold=True)
    add_bullets(
        slide,
        Inches(1.00),
        Inches(5.00),
        Inches(6.3),
        Inches(0.75),
        [
            "Uniform throttling: fixed answer rate.",
            "Targeted novelty policy: answer low-novelty items, defer high-novelty items.",
            "Capacity-aware baseline: answer more when recent overload is high.",
            "Novelty+capacity baseline: combine composition and congestion feedback.",
        ],
        size=15,
        bullet_color=ACCENT,
    )
    add_panel(slide, Inches(9.20), Inches(2.05), Inches(3.70), Inches(4.35))
    add_text(slide, Inches(9.43), Inches(2.25), Inches(2.7), Inches(0.25), "Experiment design", size=20, color=ACCENT, bold=True)
    add_text(slide, Inches(9.43), Inches(2.65), Inches(3.0), Inches(1.35), "40-round deterministic\nsimulator\n30 novelty seeds per\ncomparison\n3 regimes: stable, mid,\ncollapse", size=18, color=TEXT, bold=True)
    add_text(slide, Inches(9.43), Inches(4.55), Inches(2.8), Inches(0.25), "New checks we added", size=20, bold=True)
    add_bullets(
        slide,
        Inches(9.38),
        Inches(4.85),
        Inches(3.0),
        Inches(0.80),
        [
            "Matched answered-share control",
            "rho=0 ablation to switch off novelty",
            "Collapse boundary in (C1, xi) plane",
        ],
        size=15,
        bullet_color=ACCENT,
    )
    add_footer(slide, "Refs: Taitler et al. (2025); project simulator in sim/model.py, sim/policies.py, sim/run_experiments.py", 5)


def results_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    slide_title(slide, "Main Results", "Novelty-aware targeting beats uniform throttling across stable, mid, and collapse regimes.")
    add_badge(slide, Inches(0.55), Inches(1.95), Inches(1.0), Inches(0.24), "Our result", ACCENT)
    add_panel(slide, Inches(0.55), Inches(2.05), Inches(3.35), Inches(4.35))
    add_text(slide, Inches(0.78), Inches(2.25), Inches(2.2), Inches(0.25), "Headline numbers", size=20, color=ACCENT, bold=True)
    add_bullets(
        slide,
        Inches(0.74),
        Inches(2.58),
        Inches(2.9),
        Inches(2.1),
        [
            "Stable regime: S_T rises from 26.15 to 36.29.",
            "Mid regime: S_T rises from 19.60 to 31.56 and overload drops from 40 to 25 rounds.",
            "Collapse regime: learning still improves from 6.68 to 9.54, but overload remains saturated.",
        ],
        size=15,
        bullet_color=ACCENT,
    )
    add_text(slide, Inches(0.78), Inches(4.95), Inches(1.5), Inches(0.20), "Interpretation", size=18, bold=True)
    add_text(slide, Inches(0.78), Inches(5.22), Inches(2.8), Inches(0.55), "Targeting helps most when the forum is still productive enough for traffic composition to matter.", size=14, color=MUTED)
    add_panel(slide, Inches(4.20), Inches(2.05), Inches(4.25), Inches(4.35))
    add_picture(slide, "sim/figs/extra/seed_robustness_ST.png", Inches(4.38), Inches(2.22), Inches(3.90), Inches(3.40))
    add_text(slide, Inches(4.35), Inches(5.98), Inches(2.5), Inches(0.18), "Final learning across 30 seeds", size=11, color=MUTED)
    add_panel(slide, Inches(8.75), Inches(2.05), Inches(4.15), Inches(4.35))
    add_picture(slide, "sim/figs/extra/seed_robustness_overload.png", Inches(8.93), Inches(2.22), Inches(3.80), Inches(3.40))
    add_text(slide, Inches(8.90), Inches(5.98), Inches(2.5), Inches(0.18), "Overload-time comparison", size=11, color=MUTED)
    add_footer(slide, "Refs: seed_robustness.csv; Taitler et al. (2025)", 6)


def validation_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    slide_title(slide, "Mechanism Validation", "We tested whether the targeted gain is real, not just an artifact of adoption or answer volume.")
    add_badge(slide, Inches(0.55), Inches(1.95), Inches(1.0), Inches(0.24), "Our result", ACCENT)
    add_panel(slide, Inches(0.55), Inches(2.05), Inches(3.35), Inches(4.35))
    add_text(slide, Inches(0.78), Inches(2.25), Inches(2.5), Inches(0.25), "Two strongest checks", size=20, color=ACCENT, bold=True)
    add_bullets(
        slide,
        Inches(0.74),
        Inches(2.58),
        Inches(2.95),
        Inches(2.6),
        [
            "Matched answered share: both policies answer about 0.178 of users, but targeted still reaches S_T = 31.35 vs 20.58 for uniform.",
            "rho ablation: when novelty has no effect (rho = 0), the targeted advantage becomes exactly zero.",
            "Conclusion: the gain comes from sending better traffic to the forum, not just changing total answered volume.",
        ],
        size=15,
        bullet_color=ACCENT,
    )
    add_panel(slide, Inches(4.20), Inches(2.05), Inches(4.25), Inches(4.35))
    add_picture(slide, "sim/figs/extra/matched_answered_compare_ST.png", Inches(4.38), Inches(2.22), Inches(3.90), Inches(3.40))
    add_text(slide, Inches(4.35), Inches(5.98), Inches(2.7), Inches(0.18), "Matched answered-share control", size=11, color=MUTED)
    add_panel(slide, Inches(8.75), Inches(2.05), Inches(4.15), Inches(4.35))
    add_picture(slide, "sim/figs/extra/rho_ablation_gap_ST.png", Inches(8.93), Inches(2.22), Inches(3.80), Inches(3.40))
    add_text(slide, Inches(8.90), Inches(5.98), Inches(1.8), Inches(0.18), "rho ablation", size=11, color=MUTED)
    add_footer(slide, "Refs: matched_answered_compare.csv; rho_ablation.csv", 7)


def baselines_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    slide_title(slide, "Stronger Baselines And Tipping", "Congestion feedback helps, but novelty+capacity performs best in the mid regime.")
    add_badge(slide, Inches(0.55), Inches(1.95), Inches(1.0), Inches(0.24), "Our result", ACCENT)
    add_panel(slide, Inches(0.55), Inches(2.05), Inches(3.35), Inches(4.35))
    add_text(slide, Inches(0.78), Inches(2.25), Inches(2.8), Inches(0.25), "Mid-regime policy ranking", size=20, color=ACCENT, bold=True)
    add_bullets(
        slide,
        Inches(0.74),
        Inches(2.58),
        Inches(2.9),
        Inches(1.9),
        [
            "Uniform: S_T = 19.60, overload = 40.",
            "Capacity-aware: S_T = 22.52, overload = 25.",
            "Targeted novelty: S_T = 31.56, overload = 25.",
            "Novelty+capacity: S_T = 33.23, overload = 18.",
        ],
        size=15,
        bullet_color=ACCENT,
    )
    add_text(slide, Inches(0.78), Inches(4.95), Inches(1.5), Inches(0.20), "Phase transition", size=18, bold=True)
    add_text(slide, Inches(0.78), Inches(5.22), Inches(2.8), Inches(0.60), "First collapse points appear near C1 ~ 0.364 and xi ~ 0.074. Higher overload damage requires higher initial forum capacity.", size=14, color=MUTED)
    add_panel(slide, Inches(4.20), Inches(2.05), Inches(4.15), Inches(4.35))
    add_picture(slide, "sim/figs/extra/policy_compare_extended_ST.png", Inches(4.38), Inches(2.22), Inches(3.80), Inches(3.40))
    add_text(slide, Inches(4.35), Inches(5.98), Inches(2.5), Inches(0.18), "Extended policy comparison", size=11, color=MUTED)
    add_panel(slide, Inches(8.65), Inches(2.05), Inches(4.25), Inches(4.35))
    add_picture(slide, "sim/figs/extra/collapse_boundary_c1_xi.png", Inches(8.83), Inches(2.22), Inches(3.90), Inches(3.40))
    add_text(slide, Inches(8.80), Inches(5.98), Inches(2.3), Inches(0.18), "Explicit collapse boundary", size=11, color=MUTED)
    add_footer(slide, "Refs: policy_compare_extended.csv; collapse_boundary_c1_xi.csv", 8)


def summary_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    slide_title(slide, "Summary And Future Work", "References are listed on this slide; our contribution remains distinct from the base paper.")
    add_kicker(slide, Inches(0.55), Inches(1.95), "Close", BLUE, width=0.75)
    add_panel(slide, Inches(0.55), Inches(2.00), Inches(4.15), Inches(4.65))
    add_text(slide, Inches(0.78), Inches(2.20), Inches(2.0), Inches(0.25), "Summary", size=20, color=BLUE, bold=True)
    add_bullets(
        slide,
        Inches(0.74),
        Inches(2.50),
        Inches(3.5),
        Inches(3.40),
        [
            "Base paper insight: answering less can improve long-run learning.",
            "Our extension: forums are congestible, so routing must consider both quantity and quality of deferred traffic.",
            "Main result: novelty-aware routing dominates uniform throttling and survives strict controls.",
            "Best baseline in our mid regime combines novelty with congestion feedback.",
        ],
        size=16,
        bullet_color=BLUE,
    )
    add_panel(slide, Inches(5.05), Inches(2.00), Inches(3.85), Inches(4.65), fill=LIGHT_ORANGE, line=ACCENT)
    add_text(slide, Inches(5.28), Inches(2.20), Inches(2.0), Inches(0.25), "Future work", size=20, color=ACCENT, bold=True)
    add_bullets(
        slide,
        Inches(5.24),
        Inches(2.50),
        Inches(3.2),
        Inches(3.40),
        [
            "Estimate parameters from real forum data instead of reduced-form calibration.",
            "Replace hand-designed rules with learned routing policies.",
            "Model heterogeneous experts, delayed retraining, and multiple forums.",
            "Optimize jointly for learning, welfare, and revenue instead of one metric.",
        ],
        size=16,
        bullet_color=ACCENT,
    )
    add_panel(slide, Inches(9.25), Inches(2.00), Inches(3.65), Inches(4.65), fill=LIGHT_GREEN, line=GREEN)
    add_text(slide, Inches(9.48), Inches(2.20), Inches(2.8), Inches(0.45), "References used in the\ntalk", size=20, color=GREEN, bold=True)
    refs = "\n".join(
        [
            "- Taitler et al. (2025)",
            "- Chow (1970)",
            "- El-Yaniv (2010)",
            "- Madras et al. (2018)",
            "- Mozannar and Sontag (2020)",
            "- Hardin (1968)",
            "- Ostrom (1990)",
            "- Vickrey (1969)",
            "- Kleinrock (1975)",
        ]
    )
    add_text(slide, Inches(9.50), Inches(2.80), Inches(2.9), Inches(3.3), refs, size=16)
    add_panel(slide, Inches(0.55), Inches(6.90), Inches(11.90), Inches(0.40))
    add_text(slide, Inches(0.78), Inches(7.02), Inches(11.4), Inches(0.20), "Practice target for 8 minutes: ~45-55 seconds per slide. Keep slide 4 explicitly labeled as prior work and slides 5-8 as our contribution.", size=12, color=MUTED)
    add_footer(slide, "Refs: full bibliography in refs.bib; slide cites shown inline throughout the deck", 9)


OUT_DIR.mkdir(exist_ok=True)
title_slide()
problem_slide()
literature_slide()
base_paper_slide()
method_slide()
results_slide()
validation_slide()
baselines_slide()
summary_slide()
prs.save(PPTX_OUT)
print(PPTX_OUT)
