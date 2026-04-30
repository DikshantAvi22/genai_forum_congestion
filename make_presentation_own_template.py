from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "presentation"
PPTX_OUT = OUT_DIR / "Selective_Response_Forum_Congestion_CustomTemplate.pptx"

NAVY = RGBColor(19, 33, 68)
SLATE = RGBColor(88, 101, 123)
TEAL = RGBColor(18, 123, 137)
ORANGE = RGBColor(225, 122, 29)
RED = RGBColor(186, 56, 56)
GREEN = RGBColor(42, 117, 79)
CREAM = RGBColor(251, 248, 242)
WHITE = RGBColor(255, 255, 255)
PALE_BLUE = RGBColor(233, 242, 248)
PALE_ORANGE = RGBColor(255, 241, 227)
PALE_GREEN = RGBColor(235, 246, 238)
PALE_RED = RGBColor(252, 237, 237)
LINE = RGBColor(210, 219, 232)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.2)
    return shp


def round_rect(slide, x, y, w, h, fill, line=None, r=0.12):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = r
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.2)
    return shp


def add_text(slide, x, y, w, h, text, size=20, color=NAVY, bold=False, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, size=17, color=NAVY):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def add_pic(slide, rel, x, y, w, h):
    slide.shapes.add_picture(str(ROOT / rel), x, y, width=w, height=h)


def base_bg(slide, accent=TEAL):
    rect(slide, 0, 0, prs.slide_width, prs.slide_height, CREAM)
    rect(slide, Inches(0.45), Inches(0.38), Inches(12.43), Inches(0.13), NAVY)
    rect(slide, Inches(0.45), Inches(6.95), Inches(12.43), Inches(0.02), LINE)
    rect(slide, Inches(0.45), Inches(0.55), Inches(1.65), Inches(0.10), accent)


def title_block(slide, title, subtitle, tag, tag_fill):
    add_text(slide, Inches(0.62), Inches(0.82), Inches(8.8), Inches(0.60), title, size=28, bold=True)
    add_text(slide, Inches(0.64), Inches(1.42), Inches(9.2), Inches(0.28), subtitle, size=15, color=SLATE)
    round_rect(slide, Inches(0.62), Inches(1.78), Inches(1.15), Inches(0.26), tag_fill, None, 0.25)
    add_text(slide, Inches(0.72), Inches(1.83), Inches(0.95), Inches(0.14), tag, size=11, color=WHITE, bold=True)


def footer(slide, refs, n):
    add_text(slide, Inches(0.55), Inches(7.00), Inches(9.8), Inches(0.18), refs, size=9, color=SLATE)
    round_rect(slide, Inches(12.45), Inches(6.88), Inches(0.38), Inches(0.28), PALE_BLUE, LINE, 0.20)
    add_text(slide, Inches(12.57), Inches(6.94), Inches(0.14), Inches(0.12), str(n), size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def slide_1_title():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    rect(slide, Inches(0.6), Inches(0.55), Inches(0.18), Inches(5.9), TEAL)
    add_text(slide, Inches(1.0), Inches(0.9), Inches(6.7), Inches(1.9), "Selective Response for\nGenAI Under Forum\nCongestion", size=30, color=WHITE, bold=True)
    add_text(slide, Inches(1.02), Inches(3.0), Inches(5.7), Inches(0.3), "Game Theory 2026 mini project", size=18, color=RGBColor(220, 229, 240))
    round_rect(slide, Inches(1.02), Inches(3.55), Inches(2.35), Inches(0.34), TEAL, None, 0.18)
    add_text(slide, Inches(1.16), Inches(3.63), Inches(2.0), Inches(0.16), "Base paper + our extension", size=12, color=WHITE, bold=True)
    add_text(slide, Inches(1.02), Inches(4.15), Inches(5.8), Inches(0.3), "Dikshant Gupta (24108)  |  Paras Raina (24170)", size=18, color=WHITE)
    add_text(slide, Inches(1.02), Inches(4.52), Inches(5.2), Inches(0.24), "8-minute plan: Dikshant slides 1-5, Paras slides 6-11", size=12, color=RGBColor(214, 224, 238))

    round_rect(slide, Inches(7.4), Inches(0.95), Inches(5.0), Inches(2.15), WHITE, None, 0.08)
    add_pic(slide, "sim/figs/extra/seed_robustness_ST.png", Inches(7.62), Inches(1.12), Inches(4.56), Inches(1.62))
    round_rect(slide, Inches(7.4), Inches(3.35), Inches(5.0), Inches(2.15), WHITE, None, 0.08)
    add_pic(slide, "sim/figs/extra/collapse_boundary_c1_xi.png", Inches(7.62), Inches(3.52), Inches(4.56), Inches(1.62))
    add_text(slide, Inches(7.48), Inches(5.78), Inches(3.8), Inches(0.24), "Selective response helps, but forum congestion changes the policy design problem.", size=16, color=WHITE)
    add_text(slide, Inches(7.48), Inches(6.18), Inches(4.3), Inches(0.24), "Main message: traffic composition matters, not just answer volume.", size=16, color=RGBColor(220, 229, 240), bold=True)


def slide_2_problem():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_bg(slide, TEAL)
    title_block(slide, "Problem Setup", "When should GenAI answer directly, and when should it defer users to a human forum?", "Problem", TEAL)
    round_rect(slide, Inches(0.6), Inches(2.15), Inches(4.1), Inches(3.95), WHITE, LINE, 0.10)
    add_text(slide, Inches(0.85), Inches(2.35), Inches(2.6), Inches(0.25), "Tension", size=20, bold=True, color=TEAL)
    add_bullets(
        slide,
        Inches(0.82),
        Inches(2.70),
        Inches(3.55),
        Inches(2.65),
        [
            "Answer too often: users stop creating human discussions that could become future training signal.",
            "Defer too often: the forum gets overloaded, answers worsen, and experts burn out.",
            "Good policy must trade off immediate service against long-run ecosystem health.",
        ],
        size=18,
    )
    round_rect(slide, Inches(4.95), Inches(2.15), Inches(7.8), Inches(3.95), PALE_BLUE, LINE, 0.10)
    add_text(slide, Inches(5.20), Inches(2.35), Inches(2.8), Inches(0.25), "System view", size=20, bold=True, color=TEAL)
    round_rect(slide, Inches(5.35), Inches(3.05), Inches(1.35), Inches(0.78), WHITE, TEAL, 0.18)
    round_rect(slide, Inches(7.15), Inches(3.05), Inches(1.35), Inches(0.78), WHITE, ORANGE, 0.18)
    round_rect(slide, Inches(8.95), Inches(3.05), Inches(1.50), Inches(0.78), WHITE, GREEN, 0.18)
    round_rect(slide, Inches(10.90), Inches(3.05), Inches(1.30), Inches(0.78), WHITE, RED, 0.18)
    add_text(slide, Inches(5.70), Inches(3.31), Inches(0.7), Inches(0.16), "Users", size=16, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.43), Inches(3.31), Inches(0.8), Inches(0.16), "GenAI", size=16, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(9.28), Inches(3.31), Inches(0.9), Inches(0.16), "Forum", size=16, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(11.23), Inches(3.31), Inches(0.6), Inches(0.16), "C_t", size=16, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(6.72), Inches(3.20), Inches(0.28), Inches(0.20), "->", size=22, color=TEAL, bold=True)
    add_text(slide, Inches(8.52), Inches(3.20), Inches(0.28), Inches(0.20), "->", size=22, color=ORANGE, bold=True)
    add_text(slide, Inches(10.55), Inches(3.20), Inches(0.28), Inches(0.20), "->", size=22, color=GREEN, bold=True)
    add_text(slide, Inches(8.65), Inches(4.35), Inches(2.7), Inches(0.22), "Overload damages future capacity", size=15, color=RED, bold=True)
    add_text(slide, Inches(5.25), Inches(5.25), Inches(6.8), Inches(0.45), "Core idea: the forum is not just a knowledge source. It is a limited shared resource.", size=18, color=NAVY, bold=True)
    footer(slide, "Refs: Taitler et al. (2025); Hardin (1968); Ostrom (1990)", 2)


def slide_3_importance():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_bg(slide, ORANGE)
    title_block(slide, "Why This Matters", "The same routing problem appears in multiple real settings where human discussion creates reusable value.", "Examples", ORANGE)
    cards = [
        (Inches(0.65), "Course forums", "If GenAI answers every assignment question, fewer student-human explanations remain for later batches.", PALE_ORANGE),
        (Inches(4.48), "Developer support", "Redirected traffic can overload technical forums, making expert time the real bottleneck.", PALE_BLUE),
        (Inches(8.31), "Product communities", "Automation can solve easy issues now while silently shrinking the future knowledge base.", PALE_GREEN),
    ]
    for x, title, body, fill in cards:
        round_rect(slide, x, Inches(2.15), Inches(3.55), Inches(3.45), fill, LINE, 0.10)
        round_rect(slide, x + Inches(0.22), Inches(2.42), Inches(0.48), Inches(0.48), ORANGE if fill == PALE_ORANGE else TEAL if fill == PALE_BLUE else GREEN, None, 0.28)
        add_text(slide, x + Inches(0.82), Inches(2.45), Inches(2.2), Inches(0.22), title, size=19, bold=True)
        add_text(slide, x + Inches(0.25), Inches(3.05), Inches(3.0), Inches(1.35), body, size=18, color=SLATE)
    round_rect(slide, Inches(0.65), Inches(5.85), Inches(11.2), Inches(0.62), WHITE, LINE, 0.10)
    add_text(slide, Inches(0.92), Inches(6.03), Inches(10.7), Inches(0.22), "Takeaway: abstention is not only an accuracy decision. It is a platform-design decision.", size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, "Refs: Horvitz (1999); Amershi et al. (2014); Bansal et al. (2021)", 3)


def slide_4_lit():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_bg(slide, TEAL)
    title_block(slide, "Literature Review", "Our project lies where three literatures meet.", "Related work", TEAL)
    cols = [
        (Inches(0.65), Inches(2.15), PALE_BLUE, TEAL, "Selective response", ["Chow (1970)", "El-Yaniv (2010)", "Geifman and El-Yaniv (2017)"], "Question: when should the model abstain?"),
        (Inches(4.48), Inches(2.15), PALE_ORANGE, ORANGE, "Learning to defer", ["Madras et al. (2018)", "Mozannar and Sontag (2020)", "Kleinberg et al. (2018)"], "Question: when is human expertise more valuable than model output?"),
        (Inches(8.31), Inches(2.15), PALE_GREEN, GREEN, "Congestion / commons", ["Hardin (1968)", "Ostrom (1990)", "Kleinrock (1975)", "Vickrey (1969)"], "Question: what happens when shared systems are overloaded?"),
    ]
    for x, y, fill, accent, title, refs, line in cols:
        round_rect(slide, x, y, Inches(3.55), Inches(3.70), fill, LINE, 0.10)
        add_text(slide, x + Inches(0.22), y + Inches(0.24), Inches(2.7), Inches(0.25), title, size=19, bold=True, color=accent)
        add_text(slide, x + Inches(0.22), y + Inches(0.62), Inches(2.9), Inches(0.65), "\n".join(refs), size=16, color=NAVY)
        add_text(slide, x + Inches(0.22), y + Inches(2.08), Inches(3.0), Inches(0.85), line, size=17, color=SLATE)
    round_rect(slide, Inches(1.55), Inches(6.05), Inches(10.25), Inches(0.42), WHITE, LINE, 0.10)
    add_text(slide, Inches(1.78), Inches(6.17), Inches(9.8), Inches(0.18), "Gap: prior selective-response work does not model a congestible human forum with endogenous capacity.", size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    footer(slide, "Refs: Chow (1970); El-Yaniv (2010); Madras et al. (2018); Mozannar and Sontag (2020); Hardin (1968); Ostrom (1990)", 4)


def slide_5_basepaper():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_bg(slide, ORANGE)
    title_block(slide, "Base Paper", "This is prior work we build on, not our contribution.", "Prior work", ORANGE)
    round_rect(slide, Inches(0.65), Inches(2.10), Inches(4.05), Inches(4.25), WHITE, LINE, 0.10)
    add_text(slide, Inches(0.90), Inches(2.35), Inches(2.7), Inches(0.25), "Base paper question", size=20, bold=True, color=ORANGE)
    add_bullets(
        slide,
        Inches(0.88),
        Inches(2.72),
        Inches(3.45),
        Inches(2.30),
        [
            "Should GenAI selectively defer some users so human communities continue creating new knowledge?",
            "Main idea: less answering can improve long-run learning.",
            "Reason: deferred interactions generate future human signal.",
        ],
        size=18,
    )
    round_rect(slide, Inches(5.00), Inches(2.10), Inches(3.10), Inches(4.25), PALE_ORANGE, LINE, 0.10)
    add_text(slide, Inches(5.25), Inches(2.35), Inches(2.0), Inches(0.25), "Base paper result", size=20, bold=True, color=ORANGE)
    round_rect(slide, Inches(5.52), Inches(3.00), Inches(2.05), Inches(0.70), WHITE, ORANGE, 0.20)
    round_rect(slide, Inches(5.52), Inches(4.55), Inches(2.05), Inches(0.70), WHITE, GREEN, 0.20)
    add_text(slide, Inches(5.93), Inches(3.23), Inches(1.3), Inches(0.18), "Always answer", size=16, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.93), Inches(4.78), Inches(1.3), Inches(0.18), "Selective response", size=16, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(6.18), Inches(3.82), Inches(0.7), Inches(0.20), "vs", size=18, color=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.35), Inches(5.55), Inches(2.4), Inches(0.40), "Some deferral can beat always-answering in long-run learning.", size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    round_rect(slide, Inches(8.40), Inches(2.10), Inches(4.45), Inches(4.25), PALE_RED, LINE, 0.10)
    add_text(slide, Inches(8.65), Inches(2.35), Inches(2.8), Inches(0.50), "Limitation relevant to our project", size=20, bold=True, color=RED)
    add_bullets(
        slide,
        Inches(8.63),
        Inches(2.92),
        Inches(3.75),
        Inches(2.60),
        [
            "Forum is treated as a source of knowledge, not as a congestible shared system.",
            "No endogenous capacity decline or collapse mechanism.",
            "No explicit separation between answer volume and traffic composition effects.",
        ],
        size=18,
    )
    footer(slide, "Refs: Taitler et al. (2025)", 5)


def slide_6_ourcontrib():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_bg(slide, TEAL)
    title_block(slide, "Our Contribution", "We turn the base paper into a congestion-aware ecosystem model.", "Our work", TEAL)
    round_rect(slide, Inches(0.65), Inches(2.12), Inches(7.0), Inches(4.15), WHITE, LINE, 0.10)
    add_text(slide, Inches(0.92), Inches(2.35), Inches(2.8), Inches(0.25), "Methodology", size=20, bold=True, color=TEAL)
    add_bullets(
        slide,
        Inches(0.88),
        Inches(2.74),
        Inches(3.2),
        Inches(2.20),
        [
            "Deterministic simulator with horizon T = 40.",
            "State tracks knowledge stock, forum capacity, choice share, and answered share.",
            "Forum throughput falls under overload; capacity can decline over time.",
            "Policies compared: uniform, targeted novelty, capacity-aware, novelty+capacity.",
        ],
        size=17,
    )
    round_rect(slide, Inches(4.20), Inches(2.85), Inches(1.15), Inches(0.72), PALE_BLUE, TEAL, 0.18)
    round_rect(slide, Inches(5.75), Inches(2.85), Inches(1.15), Inches(0.72), PALE_ORANGE, ORANGE, 0.18)
    round_rect(slide, Inches(7.30), Inches(2.85), Inches(1.25), Inches(0.72), PALE_GREEN, GREEN, 0.18)
    round_rect(slide, Inches(8.95), Inches(2.85), Inches(1.00), Inches(0.72), PALE_RED, RED, 0.18)
    add_text(slide, Inches(4.48), Inches(3.10), Inches(0.6), Inches(0.16), "Users", size=16, bold=True)
    add_text(slide, Inches(6.00), Inches(3.10), Inches(0.6), Inches(0.16), "GenAI", size=16, bold=True)
    add_text(slide, Inches(7.63), Inches(3.10), Inches(0.6), Inches(0.16), "Forum", size=16, bold=True)
    add_text(slide, Inches(9.23), Inches(3.10), Inches(0.4), Inches(0.16), "C_t", size=16, bold=True)
    add_text(slide, Inches(5.35), Inches(3.04), Inches(0.24), Inches(0.18), "->", size=22, color=TEAL, bold=True)
    add_text(slide, Inches(6.90), Inches(3.04), Inches(0.24), Inches(0.18), "->", size=22, color=ORANGE, bold=True)
    add_text(slide, Inches(8.55), Inches(3.04), Inches(0.24), Inches(0.18), "->", size=22, color=GREEN, bold=True)
    add_text(slide, Inches(6.35), Inches(4.20), Inches(3.2), Inches(0.22), "Overload weakens future capacity", size=15, color=RED, bold=True)
    round_rect(slide, Inches(8.00), Inches(2.12), Inches(4.85), Inches(4.15), PALE_BLUE, LINE, 0.10)
    add_text(slide, Inches(8.28), Inches(2.35), Inches(2.2), Inches(0.25), "New checks we added", size=20, bold=True, color=TEAL)
    add_bullets(
        slide,
        Inches(8.25),
        Inches(2.76),
        Inches(4.0),
        Inches(2.25),
        [
            "Matched answered-share control",
            "rho = 0 ablation to switch off novelty",
            "Collapse boundary in (C1, xi) space",
            "30-seed robustness comparisons",
        ],
        size=18,
    )
    round_rect(slide, Inches(8.28), Inches(5.20), Inches(4.0), Inches(0.65), WHITE, LINE, 0.10)
    add_text(slide, Inches(8.48), Inches(5.42), Inches(3.6), Inches(0.20), "This is the part that is original to our team.", size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, "Refs: project simulator outputs in sim/results and sim/figs", 6)


def slide_7_replication():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_bg(slide, ORANGE)
    title_block(slide, "Replication + Extension Result", "First, we replicate the base intuition. Then we show what changes once congestion is modeled.", "Replication", ORANGE)
    round_rect(slide, Inches(0.65), Inches(2.10), Inches(3.15), Inches(4.25), WHITE, LINE, 0.10)
    add_text(slide, Inches(0.92), Inches(2.35), Inches(2.5), Inches(0.25), "Replication statement", size=20, bold=True, color=ORANGE)
    add_bullets(
        slide,
        Inches(0.88),
        Inches(2.75),
        Inches(2.6),
        Inches(2.15),
        [
            "We recover the main base-paper intuition: selective response can improve long-run learning.",
            "But once congestion is introduced, the correct question is not only how much to defer, but which traffic to defer.",
        ],
        size=18,
    )
    add_text(slide, Inches(0.92), Inches(5.35), Inches(2.5), Inches(0.18), "This slide transitions from prior work to our work.", size=13, color=SLATE)
    round_rect(slide, Inches(4.05), Inches(2.10), Inches(4.15), Inches(4.25), PALE_BLUE, LINE, 0.10)
    add_pic(slide, "sim/figs/extra/seed_robustness_ST.png", Inches(4.23), Inches(2.28), Inches(3.80), Inches(3.35))
    add_text(slide, Inches(4.20), Inches(5.88), Inches(2.8), Inches(0.18), "Learning gains across regimes", size=11, color=SLATE)
    round_rect(slide, Inches(8.45), Inches(2.10), Inches(4.40), Inches(4.25), PALE_ORANGE, LINE, 0.10)
    add_text(slide, Inches(8.72), Inches(2.35), Inches(3.2), Inches(0.25), "Numbers to say out loud", size=20, bold=True, color=ORANGE)
    add_bullets(
        slide,
        Inches(8.68),
        Inches(2.75),
        Inches(3.6),
        Inches(2.25),
        [
            "Stable: S_T 26.15 -> 36.29",
            "Mid: S_T 19.60 -> 31.56, overload 40 -> 25",
            "Collapse: S_T 6.68 -> 9.54, but overload remains 40",
        ],
        size=18,
    )
    footer(slide, "Refs: seed_robustness.csv", 7)


def slide_8_mainresults():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_bg(slide, TEAL)
    title_block(slide, "Main Results", "Novelty-aware routing improves learning because it changes the composition of forum traffic.", "Results", TEAL)
    round_rect(slide, Inches(0.65), Inches(2.12), Inches(4.00), Inches(4.20), WHITE, LINE, 0.10)
    add_text(slide, Inches(0.92), Inches(2.38), Inches(2.4), Inches(0.25), "Interpretation", size=20, bold=True, color=TEAL)
    add_bullets(
        slide,
        Inches(0.88),
        Inches(2.78),
        Inches(3.35),
        Inches(2.80),
        [
            "Uniform throttling controls only volume.",
            "Targeted policy controls composition by sending higher-value traffic to the forum.",
            "Benefit is strongest when the forum is still healthy enough to produce useful knowledge.",
        ],
        size=18,
    )
    round_rect(slide, Inches(4.95), Inches(2.12), Inches(3.80), Inches(4.20), PALE_BLUE, LINE, 0.10)
    add_pic(slide, "sim/figs/extra/seed_robustness_overload.png", Inches(5.12), Inches(2.30), Inches(3.45), Inches(3.25))
    add_text(slide, Inches(5.10), Inches(5.85), Inches(2.6), Inches(0.18), "Overload-time comparison", size=11, color=SLATE)
    round_rect(slide, Inches(9.00), Inches(2.12), Inches(3.85), Inches(4.20), PALE_GREEN, LINE, 0.10)
    add_pic(slide, "sim/figs/extra/policy_compare_extended_ST.png", Inches(9.18), Inches(2.30), Inches(3.50), Inches(3.25))
    add_text(slide, Inches(9.16), Inches(5.85), Inches(2.8), Inches(0.18), "Extended policy comparison", size=11, color=SLATE)
    footer(slide, "Refs: seed_robustness.csv; policy_compare_extended.csv", 8)


def slide_9_mechanism():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_bg(slide, RED)
    title_block(slide, "Mechanism Checks", "We tested whether the gain is real or just a reporting artifact.", "Validation", RED)
    round_rect(slide, Inches(0.65), Inches(2.12), Inches(3.35), Inches(4.20), PALE_RED, LINE, 0.10)
    add_text(slide, Inches(0.92), Inches(2.35), Inches(2.4), Inches(0.25), "Why these checks matter", size=20, bold=True, color=RED)
    add_bullets(
        slide,
        Inches(0.88),
        Inches(2.72),
        Inches(2.8),
        Inches(2.60),
        [
            "Matching raw answer rate is not enough once user choice shifts endogenously.",
            "If novelty is the mechanism, the advantage should disappear when novelty has no value.",
        ],
        size=18,
    )
    round_rect(slide, Inches(4.25), Inches(2.12), Inches(4.00), Inches(4.20), WHITE, LINE, 0.10)
    add_pic(slide, "sim/figs/extra/matched_answered_compare_ST.png", Inches(4.42), Inches(2.30), Inches(3.65), Inches(3.25))
    add_text(slide, Inches(4.40), Inches(5.84), Inches(2.8), Inches(0.18), "Matched answered-share control", size=11, color=SLATE)
    round_rect(slide, Inches(8.55), Inches(2.12), Inches(4.30), Inches(4.20), WHITE, LINE, 0.10)
    add_pic(slide, "sim/figs/extra/rho_ablation_gap_ST.png", Inches(8.72), Inches(2.30), Inches(3.95), Inches(3.25))
    add_text(slide, Inches(8.70), Inches(5.84), Inches(1.5), Inches(0.18), "rho ablation", size=11, color=SLATE)
    footer(slide, "Refs: matched_answered_compare.csv; rho_ablation.csv", 9)


def slide_10_tipping():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_bg(slide, GREEN)
    title_block(slide, "Tipping And Stronger Baselines", "Congestion-aware design matters because forum collapse appears in a structured parameter region.", "Phase view", GREEN)
    round_rect(slide, Inches(0.65), Inches(2.12), Inches(3.45), Inches(4.20), PALE_GREEN, LINE, 0.10)
    add_text(slide, Inches(0.92), Inches(2.35), Inches(2.6), Inches(0.25), "Mid-regime ranking", size=20, bold=True, color=GREEN)
    add_bullets(
        slide,
        Inches(0.88),
        Inches(2.72),
        Inches(2.9),
        Inches(2.50),
        [
            "Uniform: S_T = 19.60",
            "Capacity-aware: S_T = 22.52",
            "Targeted novelty: S_T = 31.56",
            "Novelty+capacity: S_T = 33.23",
        ],
        size=18,
    )
    add_text(slide, Inches(0.92), Inches(5.45), Inches(2.6), Inches(0.18), "Best mid-regime policy combines novelty and congestion response.", size=14, color=SLATE)
    round_rect(slide, Inches(4.35), Inches(2.12), Inches(4.10), Inches(4.20), WHITE, LINE, 0.10)
    add_pic(slide, "sim/figs/extra/collapse_boundary_c1_xi.png", Inches(4.52), Inches(2.30), Inches(3.75), Inches(3.25))
    add_text(slide, Inches(4.50), Inches(5.84), Inches(2.4), Inches(0.18), "Explicit collapse boundary", size=11, color=SLATE)
    round_rect(slide, Inches(8.75), Inches(2.12), Inches(4.10), Inches(4.20), WHITE, LINE, 0.10)
    add_pic(slide, "sim/figs/timeseries_collapse.png", Inches(8.92), Inches(2.30), Inches(3.75), Inches(3.25))
    add_text(slide, Inches(8.90), Inches(5.84), Inches(2.2), Inches(0.18), "Collapse time series", size=11, color=SLATE)
    footer(slide, "Refs: policy_compare_extended.csv; collapse_boundary_c1_xi.csv", 10)


def slide_11_summary():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_bg(slide, ORANGE)
    title_block(slide, "Summary And Future Work", "Final slide before closing: keep the distinction between prior work and our contribution explicit.", "Wrap-up", ORANGE)
    round_rect(slide, Inches(0.65), Inches(2.12), Inches(4.05), Inches(4.25), WHITE, LINE, 0.10)
    add_text(slide, Inches(0.92), Inches(2.35), Inches(2.1), Inches(0.25), "Summary", size=20, bold=True, color=ORANGE)
    add_bullets(
        slide,
        Inches(0.88),
        Inches(2.72),
        Inches(3.4),
        Inches(2.55),
        [
            "Base paper gives the selective-response intuition.",
            "Our work adds congestion, endogenous capacity, and mechanism tests.",
            "Main lesson: the right traffic should reach the forum, not just less traffic.",
        ],
        size=18,
    )
    round_rect(slide, Inches(4.95), Inches(2.12), Inches(3.90), Inches(4.25), PALE_ORANGE, LINE, 0.10)
    add_text(slide, Inches(5.22), Inches(2.35), Inches(2.0), Inches(0.25), "Future work", size=20, bold=True, color=ORANGE)
    add_bullets(
        slide,
        Inches(5.18),
        Inches(2.72),
        Inches(3.25),
        Inches(2.55),
        [
            "Calibrate parameters from real forum data.",
            "Learn routing policies instead of hand-designed rules.",
            "Model multiple forums and heterogeneous experts.",
        ],
        size=18,
    )
    round_rect(slide, Inches(9.10), Inches(2.12), Inches(3.75), Inches(4.25), PALE_BLUE, LINE, 0.10)
    add_text(slide, Inches(9.37), Inches(2.35), Inches(2.0), Inches(0.25), "References", size=20, bold=True, color=TEAL)
    add_text(
        slide,
        Inches(9.32),
        Inches(2.72),
        Inches(3.05),
        Inches(2.95),
        "Taitler et al. (2025)\nChow (1970)\nEl-Yaniv (2010)\nMadras et al. (2018)\nMozannar and Sontag (2020)\nHardin (1968)\nOstrom (1990)\nVickrey (1969)",
        size=16,
        color=NAVY,
    )
    footer(slide, "Refs: full bibliography in refs.bib", 11)


def slide_12_thanks():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, prs.slide_width, prs.slide_height, CREAM)
    rect(slide, Inches(0.55), Inches(0.48), Inches(12.2), Inches(6.35), NAVY)
    round_rect(slide, Inches(1.10), Inches(1.20), Inches(4.10), Inches(0.42), TEAL, None, 0.18)
    add_text(slide, Inches(1.30), Inches(1.32), Inches(3.5), Inches(0.18), "Selective Response for GenAI Under Forum Congestion", size=14, color=WHITE, bold=True)
    add_text(slide, Inches(1.10), Inches(2.05), Inches(5.8), Inches(1.2), "Thank You", size=34, color=WHITE, bold=True)
    add_text(slide, Inches(1.12), Inches(3.10), Inches(5.6), Inches(0.40), "Questions?", size=24, color=RGBColor(214, 224, 238), bold=True)
    add_text(slide, Inches(1.12), Inches(3.75), Inches(6.0), Inches(0.32), "Dikshant Gupta (24108)  |  Paras Raina (24170)", size=18, color=WHITE)
    round_rect(slide, Inches(7.55), Inches(1.55), Inches(4.25), Inches(3.95), WHITE, None, 0.08)
    add_pic(slide, "sim/figs/extra/policy_compare_extended_ST.png", Inches(7.80), Inches(1.82), Inches(3.75), Inches(2.95))
    add_text(slide, Inches(7.80), Inches(5.05), Inches(3.2), Inches(0.20), "Novelty + capacity baseline is strongest in the mid regime.", size=13, color=WHITE)


OUT_DIR.mkdir(exist_ok=True)
slide_1_title()
slide_2_problem()
slide_3_importance()
slide_4_lit()
slide_5_basepaper()
slide_6_ourcontrib()
slide_7_replication()
slide_8_mainresults()
slide_9_mechanism()
slide_10_tipping()
slide_11_summary()
slide_12_thanks()
prs.save(PPTX_OUT)
print(PPTX_OUT)
