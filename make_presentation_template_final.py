from __future__ import annotations

from pathlib import Path
from textwrap import dedent

import matplotlib.pyplot as plt
from matplotlib import rcParams
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "Minproject-Presentation Slide Template.pptx"
OUT = ROOT / "presentation" / "Selective_Response_Forum_Congestion_TemplateFinal.pptx"
ASSET_DIR = ROOT / "presentation" / "math_assets"

BLUE = RGBColor(68, 114, 196)
DARK = RGBColor(31, 45, 72)
RED = RGBColor(192, 0, 0)
ORANGE = RGBColor(217, 119, 6)
TEXT = RGBColor(40, 40, 40)
MUTED = RGBColor(102, 102, 102)
LINE = RGBColor(217, 217, 217)
LIGHT = RGBColor(250, 250, 252)
PALE_BLUE = RGBColor(238, 244, 253)
PALE_RED = RGBColor(253, 240, 240)
PALE_GREEN = RGBColor(240, 248, 243)
PALE_ORANGE = RGBColor(255, 245, 236)
WHITE = RGBColor(255, 255, 255)

rcParams["mathtext.fontset"] = "stix"
rcParams["font.family"] = "DejaVu Sans"


def text_box(slide, x, y, w, h, text, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullet_box(slide, x, y, w, h, items, size=18, color=TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def rounded(slide, x, y, w, h, fill, line=LINE, radius=0.12):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    return shp


def add_band(slide):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide.part.slide_layout.part.package.presentation_part.presentation.slide_width, Inches(0.22))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    return band


def footer(slide, refs, slide_no):
    text_box(slide, Inches(0.55), Inches(6.92), Inches(8.8), Inches(0.18), refs, size=9, color=MUTED)
    text_box(slide, Inches(10.65), Inches(6.9), Inches(1.7), Inches(0.18), "Game Theory 2026", size=10, color=BLUE, align=PP_ALIGN.RIGHT, bold=True)
    num = rounded(slide, Inches(12.55), Inches(6.84), Inches(0.34), Inches(0.24), WHITE, BLUE, 0.18)
    num.line.width = Pt(1.2)
    text_box(slide, Inches(12.63), Inches(6.89), Inches(0.16), Inches(0.12), str(slide_no), size=10, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


def title(slide, heading, subtitle, chip, chip_color=BLUE):
    add_band(slide)
    text_box(slide, Inches(0.6), Inches(0.48), Inches(10.3), Inches(0.36), heading, size=24, color=DARK, bold=True)
    text_box(slide, Inches(0.62), Inches(0.88), Inches(10.6), Inches(0.24), subtitle, size=13, color=MUTED)
    chip_box = rounded(slide, Inches(0.62), Inches(1.18), Inches(1.45), Inches(0.24), chip_color, chip_color, 0.22)
    chip_box.line.width = Pt(0)
    text_box(slide, Inches(0.76), Inches(1.22), Inches(1.12), Inches(0.14), chip, size=10, color=WHITE, bold=True)


def equation_png(text: str, name: str, fontsize=22) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    out = ASSET_DIR / name
    fig = plt.figure(figsize=(0.1, 0.1))
    fig.patch.set_alpha(0)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    t = ax.text(0, 0.5, f"${text}$", fontsize=fontsize, va="center", ha="left", color="black")
    fig.canvas.draw()
    bbox = t.get_window_extent(renderer=fig.canvas.get_renderer()).expanded(1.04, 1.25)
    bbox = bbox.transformed(fig.dpi_scale_trans.inverted())
    fig.savefig(out, transparent=True, bbox_inches=bbox, pad_inches=0.02, dpi=300)
    plt.close(fig)
    return out


def add_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text.strip()


def pic(slide, rel_path, x, y, w, h):
    slide.shapes.add_picture(str(ROOT / rel_path), x, y, width=w, height=h)


prs = Presentation(TEMPLATE)
slide_width = prs.slide_width
slide_height = prs.slide_height

# capture template images for reuse
title_slide_ref = prs.slides[0]
template_logo = title_slide_ref.shapes[3].image.blob
template_emblem = title_slide_ref.shapes[2].image.blob
OUT.parent.mkdir(exist_ok=True)
logo_path = OUT.parent / "template_logo.png"
emblem_path = OUT.parent / "template_emblem.jpeg"
logo_path.write_bytes(template_logo)
emblem_path.write_bytes(template_emblem)


def add_template_footer_graphics(slide):
    slide.shapes.add_picture(str(logo_path), Inches(0.25), Inches(6.63), width=Inches(0.80))
    slide.shapes.add_picture(str(emblem_path), Inches(12.08), Inches(6.18), width=Inches(0.58))


# Update title slide instead of replacing it
s = prs.slides[0]
s.shapes[1].text = "Selective Response for GenAI Under Forum Congestion"
s.shapes[4].text = "Dikshant Gupta (24108)  |  Paras Raina (24170)"
s.shapes[5].text = "Game Theory 2026"
text_box(s, Inches(0.90), Inches(1.50), Inches(8.0), Inches(0.32), "Base paper + our congestion-aware extension", size=16, color=MUTED)
text_box(s, Inches(0.90), Inches(1.84), Inches(8.2), Inches(0.32), "Story: motivation -> literature -> base paper -> our model -> results -> validation", size=14, color=MUTED)
add_notes(s, dedent("""
    Open with the question of when GenAI should answer versus defer.
    State immediately that the base paper gives the selective-response intuition, while our work adds forum congestion and capacity.
    Mention the talk structure in one sentence so the audience knows the flow.
"""))


layout = prs.slide_layouts[0]

# Slide 2
s = prs.slides.add_slide(layout)
title(s, "The Problem In One Intuition", "If AI helps too much, people stop creating the human knowledge that AI would have learned from later.", "Problem")
rounded(s, Inches(0.65), Inches(1.65), Inches(4.15), Inches(4.75), PALE_BLUE)
text_box(s, Inches(0.92), Inches(1.95), Inches(2.8), Inches(0.22), "Classroom intuition", size=20, color=BLUE, bold=True)
bullet_box(s, Inches(0.88), Inches(2.34), Inches(3.45), Inches(2.9), [
    "If AI answers every forum question, fewer useful human discussions happen.",
    "Those missing discussions reduce future training signal.",
    "But sending everyone to the forum also creates congestion."
], size=18)
text_box(s, Inches(0.92), Inches(5.50), Inches(3.45), Inches(0.40), "So the policy question is not only 'answer or not'. It is 'what kind of traffic should reach humans?'", size=16, color=DARK, bold=True)
rounded(s, Inches(5.05), Inches(1.65), Inches(7.15), Inches(4.75), WHITE)
text_box(s, Inches(5.32), Inches(1.96), Inches(2.4), Inches(0.22), "System view", size=20, color=BLUE, bold=True)
rounded(s, Inches(5.45), Inches(3.00), Inches(1.20), Inches(0.75), WHITE, BLUE)
rounded(s, Inches(7.20), Inches(3.00), Inches(1.20), Inches(0.75), WHITE, BLUE)
rounded(s, Inches(8.95), Inches(3.00), Inches(1.30), Inches(0.75), WHITE, BLUE)
rounded(s, Inches(10.85), Inches(3.00), Inches(1.10), Inches(0.75), WHITE, RED)
text_box(s, Inches(5.76), Inches(3.27), Inches(0.55), Inches(0.12), "Users", size=16, bold=True, align=PP_ALIGN.CENTER)
text_box(s, Inches(7.51), Inches(3.27), Inches(0.55), Inches(0.12), "GenAI", size=16, bold=True, align=PP_ALIGN.CENTER)
text_box(s, Inches(9.29), Inches(3.27), Inches(0.65), Inches(0.12), "Forum", size=16, bold=True, align=PP_ALIGN.CENTER)
text_box(s, Inches(11.15), Inches(3.27), Inches(0.45), Inches(0.12), "C_t", size=16, bold=True, align=PP_ALIGN.CENTER)
text_box(s, Inches(6.70), Inches(3.18), Inches(0.24), Inches(0.12), "->", size=22, color=BLUE, bold=True)
text_box(s, Inches(8.45), Inches(3.18), Inches(0.24), Inches(0.12), "->", size=22, color=BLUE, bold=True)
text_box(s, Inches(10.38), Inches(3.18), Inches(0.24), Inches(0.12), "->", size=22, color=RED, bold=True)
text_box(s, Inches(6.08), Inches(4.72), Inches(5.3), Inches(0.55), "Missing idea in the base paper: the forum is a limited-capacity shared resource.", size=20, color=DARK, bold=True, align=PP_ALIGN.CENTER)
footer(s, "Refs: Taitler et al. (2025); Hardin (1968); Ostrom (1990)", 2)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    Explain the story as a feedback loop.
    The left box gives the simplest intuition: AI can cannibalize the human forum that produces future knowledge.
    The right diagram adds the missing systems idea: the forum has capacity, so too much deferral also creates a cost.
"""))

# Slide 3
s = prs.slides.add_slide(layout)
title(s, "Why This Matters", "The same selective-response problem appears in multiple real human-AI ecosystems.", "Examples")
for i, (x, title_txt, body, fill) in enumerate([
    (0.72, "Course forums", "Student discussions can be reusable learning assets for future batches.", PALE_ORANGE),
    (4.40, "Developer forums", "Expert communities can be overloaded if AI redirects too many queries.", PALE_BLUE),
    (8.08, "Product support", "Automation can solve today's issue while weakening tomorrow's human knowledge base.", PALE_GREEN),
]):
    rounded(s, Inches(x), Inches(1.8), Inches(3.25), Inches(4.85), fill)
    text_box(s, Inches(x + 0.22), Inches(2.10), Inches(2.2), Inches(0.22), title_txt, size=20, color=BLUE if i == 1 else ORANGE if i == 0 else RGBColor(0, 120, 90), bold=True)
    text_box(s, Inches(x + 0.22), Inches(2.55), Inches(2.8), Inches(1.25), body, size=18)
    text_box(s, Inches(x + 0.22), Inches(5.30), Inches(2.8), Inches(0.50), "In all three cases, human interaction is not just a backup. It is a productive input.", size=15, color=MUTED)
footer(s, "Refs: Horvitz (1999); Amershi et al. (2014); Bansal et al. (2021)", 3)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    Give one crisp example from each box.
    The point is not that these platforms are identical.
    The point is that in each case, human interaction generates reusable value, so routing policy affects the ecosystem.
"""))

# Slide 4
s = prs.slides.add_slide(layout)
title(s, "Literature Review", "Our project sits where selective classification, learning to defer, and congestion economics meet.", "Literature")
for i, (x, heading, refs, q, fill, color) in enumerate([
    (0.72, "Selective response", "Chow (1970)\nEl-Yaniv (2010)", "When should the model abstain?", PALE_BLUE, BLUE),
    (4.40, "Learning to defer", "Madras et al. (2018)\nMozannar and Sontag (2020)", "When should humans take over?", PALE_ORANGE, ORANGE),
    (8.08, "Congestion / commons", "Hardin (1968)\nOstrom (1990)\nKleinrock (1975)", "What happens when a shared system is overloaded?", PALE_GREEN, RGBColor(0, 120, 90)),
]):
    rounded(s, Inches(x), Inches(1.85), Inches(3.25), Inches(4.6), fill)
    text_box(s, Inches(x + 0.20), Inches(2.12), Inches(2.7), Inches(0.22), heading, size=19, color=color, bold=True)
    text_box(s, Inches(x + 0.20), Inches(2.60), Inches(2.7), Inches(0.90), refs, size=17)
    text_box(s, Inches(x + 0.20), Inches(4.55), Inches(2.8), Inches(0.85), q, size=18, color=MUTED)
rounded(s, Inches(1.30), Inches(6.05), Inches(9.95), Inches(0.38), WHITE)
text_box(s, Inches(1.58), Inches(6.14), Inches(9.4), Inches(0.14), "Gap: prior selective-response work does not model the forum as a congestible, state-dependent resource.", size=16, bold=True, align=PP_ALIGN.CENTER)
footer(s, "Refs: Chow (1970); El-Yaniv (2010); Madras et al. (2018); Mozannar and Sontag (2020); Hardin (1968); Ostrom (1990)", 4)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    This slide is the bridge to the professor-facing contribution statement.
    Selective response says abstention matters.
    Learning to defer says humans can be the right expert.
    Congestion theory says shared systems degrade when overloaded.
    Our paper combines those three ideas.
"""))

# Slide 5
s = prs.slides.add_slide(layout)
title(s, "Base Paper: Idea, Result, Limitation", "This slide is explicitly prior work and should be spoken as prior work.", "Base paper", ORANGE)
rounded(s, Inches(0.7), Inches(1.82), Inches(4.0), Inches(4.7), WHITE)
text_box(s, Inches(0.96), Inches(2.10), Inches(2.7), Inches(0.20), "Base paper claim", size=20, color=ORANGE, bold=True)
bullet_box(s, Inches(0.92), Inches(2.48), Inches(3.3), Inches(2.55), [
    "Selective response can improve long-run learning.",
    "Reason: unanswered queries can create future human-generated signal.",
    "So always answering is not necessarily optimal."
], size=18)
rounded(s, Inches(5.05), Inches(1.82), Inches(2.65), Inches(4.7), PALE_ORANGE)
text_box(s, Inches(5.30), Inches(2.10), Inches(1.9), Inches(0.20), "Intuition", size=20, color=ORANGE, bold=True)
rounded(s, Inches(5.58), Inches(3.00), Inches(1.60), Inches(0.68), WHITE, ORANGE)
rounded(s, Inches(5.58), Inches(4.46), Inches(1.60), Inches(0.68), WHITE, BLUE)
text_box(s, Inches(5.90), Inches(3.24), Inches(0.95), Inches(0.12), "Always answer", size=15, bold=True, align=PP_ALIGN.CENTER)
text_box(s, Inches(5.83), Inches(4.70), Inches(1.10), Inches(0.12), "Selective response", size=15, bold=True, align=PP_ALIGN.CENTER)
text_box(s, Inches(5.15), Inches(5.55), Inches(2.0), Inches(0.42), "Sometimes less answering creates more future learning.", size=16, bold=True, align=PP_ALIGN.CENTER)
rounded(s, Inches(8.0), Inches(1.82), Inches(4.25), Inches(4.7), PALE_RED)
text_box(s, Inches(8.28), Inches(2.10), Inches(2.6), Inches(0.20), "Limitation for our setting", size=20, color=RED, bold=True)
bullet_box(s, Inches(8.24), Inches(2.48), Inches(3.55), Inches(2.55), [
    "The forum is not capacity-limited.",
    "There is no overload damage or collapse.",
    "Volume effects and composition effects are not separated."
], size=18)
footer(s, "Refs: Taitler et al. (2025)", 5)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    State clearly that this is not our contribution.
    We replicate the selective-response intuition, but the systems limitation becomes important once a real forum is involved.
"""))

# Slide 6 model flowchart
s = prs.slides.add_slide(layout)
title(s, "How The Equations Work", "Read this slide left to right: each box corresponds to one term in the model.", "Flowchart", BLUE)
rounded(s, Inches(0.72), Inches(1.82), Inches(11.55), Inches(4.85), WHITE)
text_box(s, Inches(0.98), Inches(2.06), Inches(3.3), Inches(0.20), "Equation-aware flowchart", size=20, color=BLUE, bold=True)

# top pipeline
top_y = 2.58
box_h = 0.70
eq_h = 0.44
desc_h = 0.36

steps = [
    (0.95, 1.45, "1. Routing", "Policy chooses\nanswered vs deferred", r"p_t^{\mathrm{answered}}", "What AI answers"),
    (2.90, 1.55, "2. Forum load", "Deferred users reach\nforum", r"m_t = 1 - p_t^{\mathrm{answered}}", "Load on forum"),
    (5.00, 1.55, "3. Productive forum", "Useful throughput and\nnovelty matter", r"G_t = \kappa\, q_t\, n_t", "Knowledge gain"),
    (7.10, 1.45, "4. Knowledge stock", "New knowledge is\naccumulated", r"S_{t+1}=S_t+G_t", "Future AI improves"),
    (9.05, 1.80, "5. Capacity update", "Overload changes next\nperiod forum health", r"C_{t+1}=(1-\delta)C_t+\eta R-\xi(m_t-C_t)^+", "Forum may weaken"),
]

eq_paths = [
    equation_png(r"p_t^{\mathrm{answered}}", "flow_eq_1.png", fontsize=19),
    equation_png(r"m_t = 1 - p_t^{\mathrm{answered}}", "flow_eq_2.png", fontsize=18),
    equation_png(r"G_t = \kappa\, q_t\, n_t", "flow_eq_3.png", fontsize=18),
    equation_png(r"S_{t+1}=S_t+G_t", "flow_eq_4.png", fontsize=19),
    equation_png(r"C_{t+1}=(1-\delta)C_t+\eta R-\xi(m_t-C_t)^+", "flow_eq_5.png", fontsize=16),
]

fills = [PALE_BLUE, PALE_BLUE, PALE_GREEN, PALE_GREEN, PALE_ORANGE]
for idx, ((x, w, title_txt, mid_txt, eq_txt, desc_txt), eq_path, fill) in enumerate(zip(steps, eq_paths, fills)):
    rounded(s, Inches(x), Inches(top_y), Inches(w), Inches(box_h), fill, BLUE if idx < 4 else RED)
    text_box(s, Inches(x + 0.08), Inches(top_y + 0.10), Inches(w - 0.16), Inches(0.16), title_txt, size=15, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, Inches(x + 0.10), Inches(top_y + 0.30), Inches(w - 0.20), Inches(0.22), mid_txt, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    rounded(s, Inches(x), Inches(top_y + 0.92), Inches(w), Inches(eq_h), WHITE, LINE)
    s.shapes.add_picture(str(eq_path), Inches(x + 0.08), Inches(top_y + 1.01), width=Inches(w - 0.16))
    text_box(s, Inches(x + 0.05), Inches(top_y + 1.45), Inches(w - 0.10), Inches(desc_h), desc_txt, size=12, color=MUTED, align=PP_ALIGN.CENTER)

for x in [2.50, 4.55, 6.65, 8.70]:
    text_box(s, Inches(x), Inches(2.86), Inches(0.26), Inches(0.12), "->", size=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# bottom overload loop
rounded(s, Inches(3.15), Inches(5.35), Inches(2.20), Inches(0.70), PALE_RED, RED)
rounded(s, Inches(6.10), Inches(5.35), Inches(2.20), Inches(0.70), PALE_BLUE, BLUE)
text_box(s, Inches(3.42), Inches(5.60), Inches(1.65), Inches(0.12), "If m_t > C_t", size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
text_box(s, Inches(6.35), Inches(5.60), Inches(1.70), Inches(0.12), "Next capacity falls", size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
text_box(s, Inches(5.38), Inches(5.53), Inches(0.34), Inches(0.12), "->", size=22, color=RED, bold=True, align=PP_ALIGN.CENTER)
text_box(s, Inches(8.72), Inches(5.34), Inches(2.75), Inches(0.70), "This feedback loop is the main extension over the base paper.", size=15, color=DARK, bold=True, align=PP_ALIGN.CENTER)
footer(s, "Refs: simulator structure from paper/main.tex and sim/model.py", 6)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    Use this as the story slide for the whole model.
    Read it left to right and point to one equation box at a time.
    Routing determines answered share.
    Answered share determines forum load.
    Productive forum traffic generates knowledge gain.
    Knowledge updates the stock.
    If load exceeds capacity, the next-period forum becomes weaker.
    End on the bottom overload loop because that is the core new idea.
"""))

# Slide 7 math
s = prs.slides.add_slide(layout)
title(s, "Core Equations And Notation", "To satisfy the mathematical expectation, we make the mechanism explicit but keep the notation light.", "Math", BLUE)
eq1 = equation_png(r"m_t = 1 - p_t^{\mathrm{answered}}", "eq_clean_1.png", fontsize=24)
eq2 = equation_png(r"S_{t+1} = S_t + \kappa\, q_t\, n_t", "eq_clean_2.png", fontsize=24)
eq3 = equation_png(r"C_{t+1} = \max\{0,(1-\delta)C_t + \eta R - \xi (m_t - C_t)^{+}\}", "eq_clean_3.png", fontsize=22)
rounded(s, Inches(0.7), Inches(1.85), Inches(6.85), Inches(4.55), WHITE)
text_box(s, Inches(0.96), Inches(2.08), Inches(2.2), Inches(0.20), "Model in 3 lines", size=20, color=BLUE, bold=True)
s.shapes.add_picture(str(eq1), Inches(1.02), Inches(2.55), width=Inches(3.9))
s.shapes.add_picture(str(eq2), Inches(1.02), Inches(3.35), width=Inches(3.9))
s.shapes.add_picture(str(eq3), Inches(1.02), Inches(4.22), width=Inches(5.9))
text_box(s, Inches(1.02), Inches(5.45), Inches(5.8), Inches(0.24), "Here, q_t = throughput and n_t = novelty factor; (x)^{+} means max(x,0).", size=13, color=MUTED)
rounded(s, Inches(7.85), Inches(1.85), Inches(4.4), Inches(4.55), PALE_BLUE)
text_box(s, Inches(8.12), Inches(2.08), Inches(2.0), Inches(0.20), "How to read it", size=20, color=BLUE, bold=True)
bullet_box(s, Inches(8.08), Inches(2.48), Inches(3.55), Inches(2.15), [
    "Line 1: how much traffic reaches the forum.",
    "Line 2: forum traffic creates new knowledge.",
    "Line 3: overload reduces future forum capacity."
], size=17)
rounded(s, Inches(8.12), Inches(5.00), Inches(3.9), Inches(0.95), WHITE)
text_box(s, Inches(8.30), Inches(5.20), Inches(3.55), Inches(0.40), "Intuition: the forum helps learning today, but too much load can damage learning tomorrow.", size=15, color=DARK, bold=True, align=PP_ALIGN.CENTER)
footer(s, "Refs: model equations from paper/main.tex", 7)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    Spend about 45-60 seconds here.
    Say: the slide has only three lines.
    First line tells us how much load reaches the forum.
    Second line tells us how the forum turns traffic into new knowledge.
    Third line tells us why congestion matters: overload reduces future capacity.
    Mention that q_t is throughput and n_t is novelty factor only if asked.
"""))

# Slide 8 design
s = prs.slides.add_slide(layout)
title(s, "Experiment Design", "We test both replication and our new congestion-aware claims.", "Design", BLUE)
for i, (x, heading, bullets_list, fill) in enumerate([
    (0.75, "Policies", ["Uniform throttling", "Targeted novelty policy", "Capacity-aware baseline", "Novelty + capacity baseline"], PALE_BLUE),
    (4.45, "Regimes", ["Stable: C1 = 0.85, xi = 0.04", "Mid: C1 = 0.60, xi = 0.10", "Collapse: C1 = 0.35, xi = 0.35", "30 novelty seeds"], PALE_ORANGE),
    (8.15, "Metrics", ["Final learning S_T", "Minimum capacity", "Overload time", "Average answered share", "Welfare and revenue proxies"], PALE_GREEN),
]):
    rounded(s, Inches(x), Inches(1.9), Inches(3.25), Inches(4.9), fill)
    text_box(s, Inches(x + 0.22), Inches(2.18), Inches(2.0), Inches(0.20), heading, size=20, color=BLUE if i == 0 else ORANGE if i == 1 else RGBColor(0, 120, 90), bold=True)
    bullet_box(s, Inches(x + 0.18), Inches(2.58), Inches(2.8), Inches(3.2), bullets_list, size=17)
footer(s, "Refs: sim/results/summary.md; sim/results/extra/summary_extra.md", 8)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    This slide tells the professor exactly what we compare and how.
    Mention that the simulator is deterministic apart from novelty draws.
    That is why seed variation here is meaningful but controlled.
"""))

# Slide 9 results
s = prs.slides.add_slide(layout)
title(s, "Main Results", "Novelty-aware routing dominates uniform throttling across regimes.", "Results", ORANGE)
rounded(s, Inches(0.72), Inches(1.88), Inches(3.2), Inches(4.8), WHITE)
text_box(s, Inches(0.98), Inches(2.16), Inches(2.2), Inches(0.20), "Results", size=20, color=ORANGE, bold=True)
bullet_box(s, Inches(0.94), Inches(2.55), Inches(2.6), Inches(2.7), [
    "Targeted novelty-aware deferral outperforms uniform throttling in all regimes.",
    "The largest gain appears in the mid regime: 19.60 -> 31.56, with overload falling from 40 to 25.",
    "Even in collapse, targeted routing still improves final learning: 6.68 -> 9.54."
], size=16)
text_box(s, Inches(0.98), Inches(5.60), Inches(2.55), Inches(0.42), "Interpretation: the forum receives more useful traffic, not just more traffic.", size=15, color=MUTED)
rounded(s, Inches(4.18), Inches(1.88), Inches(4.05), Inches(4.8), PALE_BLUE)
pic(s, "sim/figs/extra/seed_robustness_ST.png", Inches(4.38), Inches(2.10), Inches(3.65), Inches(3.55))
rounded(s, Inches(8.48), Inches(1.88), Inches(4.05), Inches(4.8), PALE_BLUE)
pic(s, "sim/figs/extra/seed_robustness_overload.png", Inches(8.68), Inches(2.10), Inches(3.65), Inches(3.55))
footer(s, "Refs: seed_robustness.csv", 9)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    Start with the middle regime because that is where the economics is clearest.
    Then mention stable and collapse briefly.
    The audience should leave this slide knowing that targeted novelty helps both learning and overload in the non-collapse region.
"""))

# Slide 10 validation
s = prs.slides.add_slide(layout)
title(s, "Mechanism Validation", "We test whether the targeted gain is truly a composition effect.", "Validation", RED)
rounded(s, Inches(0.72), Inches(1.88), Inches(3.0), Inches(4.8), PALE_RED)
text_box(s, Inches(0.98), Inches(2.16), Inches(2.0), Inches(0.20), "Two strongest checks", size=20, color=RED, bold=True)
bullet_box(s, Inches(0.94), Inches(2.55), Inches(2.4), Inches(2.8), [
    "Matched answered-share control",
    "rho = 0 ablation",
    "Together, they isolate the intended mechanism"
], size=18)
rounded(s, Inches(4.02), Inches(1.88), Inches(4.0), Inches(4.8), WHITE)
pic(s, "sim/figs/extra/matched_answered_compare_ST.png", Inches(4.22), Inches(2.10), Inches(3.6), Inches(3.55))
text_box(s, Inches(4.25), Inches(5.82), Inches(3.4), Inches(0.25), "Even at matched answered share, targeted reaches 31.35 vs 20.58.", size=14, color=MUTED, align=PP_ALIGN.CENTER)
rounded(s, Inches(8.32), Inches(1.88), Inches(4.2), Inches(4.8), WHITE)
pic(s, "sim/figs/extra/rho_ablation_gap_ST.png", Inches(8.52), Inches(2.10), Inches(3.8), Inches(3.55))
text_box(s, Inches(8.55), Inches(5.82), Inches(3.4), Inches(0.25), "When novelty has zero value, the targeted advantage disappears.", size=14, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, "Refs: matched_answered_compare.csv; rho_ablation.csv", 10)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    This is the professor's slide.
    Check 1 says the result is not just because one policy answers more.
    Check 2 says the result vanishes exactly when the novelty mechanism is switched off.
    Together, these make the causal story much stronger.
"""))

# Slide 11 tipping/summary
s = prs.slides.add_slide(layout)
title(s, "Tipping, Summary, Future Work", "We end with the systems insight, the main lesson, and what remains open.", "Summary", BLUE)
rounded(s, Inches(0.72), Inches(1.85), Inches(4.0), Inches(4.9), PALE_GREEN)
text_box(s, Inches(0.98), Inches(2.12), Inches(2.1), Inches(0.20), "Tipping insight", size=20, color=RGBColor(0, 120, 90), bold=True)
bullet_box(s, Inches(0.94), Inches(2.50), Inches(3.25), Inches(2.55), [
    "Low initial capacity and high overload damage create collapse.",
    "Collapse is a region, not a random accident.",
    "That is why forum health must be part of the policy objective."
], size=18)
pic(s, "sim/figs/extra/collapse_boundary_c1_xi.png", Inches(0.96), Inches(5.10), Inches(3.45), Inches(1.28))
rounded(s, Inches(5.0), Inches(1.85), Inches(3.35), Inches(4.9), PALE_BLUE)
text_box(s, Inches(5.26), Inches(2.12), Inches(2.5), Inches(0.20), "Key takeaways", size=20, color=BLUE, bold=True)
bullet_box(s, Inches(5.22), Inches(2.50), Inches(2.65), Inches(2.55), [
    "Which queries are deferred matters as much as how many are deferred.",
    "The rho = 0 ablation shows that the targeted advantage comes from novelty.",
    "The model shows a clear collapse region in the parameter space."
], size=17)
rounded(s, Inches(8.62), Inches(1.85), Inches(3.60), Inches(4.9), PALE_ORANGE)
text_box(s, Inches(8.88), Inches(2.12), Inches(2.1), Inches(0.20), "Future work", size=20, color=ORANGE, bold=True)
bullet_box(s, Inches(8.84), Inches(2.50), Inches(2.9), Inches(2.1), [
    "Calibrate with real forum data.",
    "Learn policies rather than hand-design them.",
    "Model multiple forums and heterogeneous experts."
], size=18)
text_box(s, Inches(8.88), Inches(5.45), Inches(2.8), Inches(0.45), "References: Taitler et al. (2025); Chow (1970); El-Yaniv (2010); Madras et al. (2018); Hardin (1968); Ostrom (1990).", size=12, color=MUTED)
footer(s, "Refs: collapse_boundary_c1_xi.csv; policy_compare_extended.csv", 11)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    Use the left box to explain why tipping matters.
    Then summarize in three lines.
    Close with future work only briefly.
    Keep the final sentence simple: in congested systems, composition matters as much as volume.
"""))

# Slide 12 references
s = prs.slides.add_slide(layout)
title(s, "References", "Key sources used in the presentation and project report.", "References", BLUE)
rounded(s, Inches(0.72), Inches(1.85), Inches(5.55), Inches(4.9), PALE_BLUE)
text_box(s, Inches(0.98), Inches(2.12), Inches(2.0), Inches(0.20), "Core paper", size=20, color=BLUE, bold=True)
text_box(s, Inches(0.98), Inches(2.52), Inches(4.6), Inches(0.65), "Taitler, S. et al. (2025). Selective Response Strategies for GenAI.", size=18)
text_box(s, Inches(0.98), Inches(3.35), Inches(2.6), Inches(0.20), "Selective response / defer", size=20, color=BLUE, bold=True)
text_box(s, Inches(0.98), Inches(3.74), Inches(4.7), Inches(1.25), "Chow (1970)\nEl-Yaniv (2010)\nMadras et al. (2018)\nMozannar and Sontag (2020)", size=17)
rounded(s, Inches(6.65), Inches(1.85), Inches(5.55), Inches(4.9), PALE_GREEN)
text_box(s, Inches(6.92), Inches(2.12), Inches(2.2), Inches(0.20), "Congestion / systems", size=20, color=BLUE, bold=True)
text_box(s, Inches(6.92), Inches(2.52), Inches(4.5), Inches(1.2), "Hardin (1968)\nOstrom (1990)\nVickrey (1969)\nKleinrock (1975)", size=17)
text_box(s, Inches(6.92), Inches(4.15), Inches(2.8), Inches(0.20), "Project artifacts", size=20, color=BLUE, bold=True)
text_box(s, Inches(6.92), Inches(4.55), Inches(4.7), Inches(1.2), "paper/main.tex\nsim/results/extra/seed_robustness.csv\nsim/results/extra/matched_answered_compare.csv\nsim/results/extra/policy_compare_extended.csv", size=16)
footer(s, "Refs: full bibliography in refs.bib", 12)
add_template_footer_graphics(s)
add_notes(s, dedent("""
    Use this slide only if asked for citations or if the professor wants a quick source summary.
    Otherwise move through it quickly and then go to the thank-you slide.
"""))

# Thank you slide
s = prs.slides.add_slide(layout)
add_band(s)
text_box(s, Inches(0.9), Inches(1.2), Inches(5.0), Inches(0.5), "Thank You", size=32, color=DARK, bold=True)
text_box(s, Inches(0.92), Inches(2.0), Inches(2.0), Inches(0.28), "Questions?", size=24, color=BLUE, bold=True)
text_box(s, Inches(0.92), Inches(2.8), Inches(4.0), Inches(0.25), "Dikshant Gupta (24108)", size=18, color=DARK)
text_box(s, Inches(0.92), Inches(3.16), Inches(4.0), Inches(0.25), "Paras Raina (24170)", size=18, color=DARK)
rounded(s, Inches(6.2), Inches(1.2), Inches(5.4), Inches(4.7), PALE_BLUE)
pic(s, "sim/figs/extra/policy_compare_extended_ST.png", Inches(6.45), Inches(1.48), Inches(4.9), Inches(3.65))
footer(s, " ", 13)
add_template_footer_graphics(s)
add_notes(s, "Pause here for questions. If asked for the main contribution in one line: we extend selective response to a congested human-forum setting and show that traffic composition matters.")

prs.save(OUT)
print(OUT)
