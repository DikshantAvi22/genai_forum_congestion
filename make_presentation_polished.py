from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "presentation" / "Selective_Response_Forum_Congestion_Polished.pptx"

BG = RGBColor(248, 246, 240)
NAVY = RGBColor(24, 38, 66)
TEAL = RGBColor(28, 126, 124)
ORANGE = RGBColor(220, 124, 48)
GREEN = RGBColor(58, 122, 87)
RED = RGBColor(184, 64, 64)
MUTED = RGBColor(100, 108, 122)
LINE = RGBColor(216, 219, 226)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(239, 242, 247)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.2)
    return shp


def rrect(slide, x, y, w, h, fill, line=None, r=0.14):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = r
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.0)
    return shp


def txt(slide, x, y, w, h, text, size=20, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, size=18, color=NAVY):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def pic(slide, rel, x, y, w, h):
    slide.shapes.add_picture(str(ROOT / rel), x, y, width=w, height=h)


def base(slide, accent=TEAL):
    rect(slide, 0, 0, prs.slide_width, prs.slide_height, BG)
    rect(slide, Inches(0.4), Inches(0.35), Inches(12.55), Inches(0.12), NAVY)
    rect(slide, Inches(0.4), Inches(6.95), Inches(12.55), Inches(0.02), LINE)
    rect(slide, Inches(0.4), Inches(0.55), Inches(1.5), Inches(0.1), accent)


def header(slide, title, subtitle, chip, chip_color):
    txt(slide, Inches(0.58), Inches(0.82), Inches(8.6), Inches(0.46), title, size=27, bold=True)
    txt(slide, Inches(0.60), Inches(1.28), Inches(9.4), Inches(0.24), subtitle, size=14, color=MUTED)
    rrect(slide, Inches(0.60), Inches(1.68), Inches(1.1), Inches(0.24), chip_color, None, 0.22)
    txt(slide, Inches(0.73), Inches(1.73), Inches(0.82), Inches(0.12), chip, size=11, color=WHITE, bold=True)


def foot(slide, refs, n):
    txt(slide, Inches(0.5), Inches(7.0), Inches(10), Inches(0.15), refs, size=9, color=MUTED)
    rrect(slide, Inches(12.48), Inches(6.88), Inches(0.34), Inches(0.26), PALE, LINE, 0.18)
    txt(slide, Inches(12.58), Inches(6.94), Inches(0.12), Inches(0.1), str(n), size=11, bold=True, align=PP_ALIGN.CENTER)


def title_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    rect(s, Inches(0.55), Inches(0.6), Inches(0.2), Inches(5.8), TEAL)
    txt(s, Inches(0.95), Inches(0.9), Inches(6.3), Inches(1.4), "Selective Response for\nGenAI Under Forum Congestion", size=29, color=WHITE, bold=True)
    txt(s, Inches(0.98), Inches(2.7), Inches(5), Inches(0.25), "Game Theory 2026 mini project", size=18, color=RGBColor(220, 226, 235))
    txt(s, Inches(0.98), Inches(3.25), Inches(6.0), Inches(0.25), "Dikshant Gupta (24108)   |   Paras Raina (24170)", size=18, color=WHITE)
    txt(s, Inches(0.98), Inches(3.62), Inches(5.4), Inches(0.18), "12 slides total: title + 10 content + thank you", size=12, color=RGBColor(220, 226, 235))
    rrect(s, Inches(7.35), Inches(1.00), Inches(5.0), Inches(2.05), WHITE, None, 0.08)
    pic(s, "sim/figs/extra/seed_robustness_ST.png", Inches(7.55), Inches(1.18), Inches(4.6), Inches(1.55))
    rrect(s, Inches(7.35), Inches(3.45), Inches(5.0), Inches(2.05), WHITE, None, 0.08)
    pic(s, "sim/figs/extra/collapse_boundary_c1_xi.png", Inches(7.55), Inches(3.63), Inches(4.6), Inches(1.55))
    txt(s, Inches(7.45), Inches(5.95), Inches(4.6), Inches(0.45), "Big picture: selective response can help learning, but forums can also collapse under load.", size=17, color=WHITE)


def s2_problem():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, TEAL)
    header(s, "The Problem In One Intuition", "GenAI and the human forum are connected. What helps now may hurt the system later.", "Problem", TEAL)
    rrect(s, Inches(0.62), Inches(2.08), Inches(4.15), Inches(4.15), WHITE, LINE, 0.08)
    txt(s, Inches(0.90), Inches(2.35), Inches(2.3), Inches(0.22), "Think of a classroom forum", size=20, bold=True, color=TEAL)
    bullets(s, Inches(0.86), Inches(2.75), Inches(3.45), Inches(2.5), [
        "If AI answers every doubt, fewer useful student discussions get created.",
        "If AI sends everyone to the forum, the forum gets crowded and less useful.",
        "So the policy question is: which queries should AI answer, and which should reach humans?"
    ], size=18)
    rrect(s, Inches(5.10), Inches(2.08), Inches(7.65), Inches(4.15), PALE, LINE, 0.08)
    rrect(s, Inches(5.45), Inches(3.00), Inches(1.25), Inches(0.72), WHITE, TEAL, 0.16)
    rrect(s, Inches(7.25), Inches(3.00), Inches(1.25), Inches(0.72), WHITE, ORANGE, 0.16)
    rrect(s, Inches(9.05), Inches(3.00), Inches(1.35), Inches(0.72), WHITE, GREEN, 0.16)
    rrect(s, Inches(10.95), Inches(3.00), Inches(1.15), Inches(0.72), WHITE, RED, 0.16)
    txt(s, Inches(5.79), Inches(3.24), Inches(0.5), Inches(0.14), "Users", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(7.58), Inches(3.24), Inches(0.55), Inches(0.14), "GenAI", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(9.40), Inches(3.24), Inches(0.6), Inches(0.14), "Forum", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(11.27), Inches(3.24), Inches(0.4), Inches(0.14), "C_t", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(6.72), Inches(3.15), Inches(0.24), Inches(0.12), "->", size=22, bold=True, color=TEAL)
    txt(s, Inches(8.52), Inches(3.15), Inches(0.24), Inches(0.12), "->", size=22, bold=True, color=ORANGE)
    txt(s, Inches(10.45), Inches(3.15), Inches(0.24), Inches(0.12), "->", size=22, bold=True, color=GREEN)
    txt(s, Inches(6.0), Inches(4.35), Inches(5.6), Inches(0.5), "Forum capacity is the missing idea: if load exceeds capacity, future human help becomes weaker.", size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    foot(s, "Refs: Taitler et al. (2025); Hardin (1968); Ostrom (1990)", 2)


def s3_examples():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, ORANGE)
    header(s, "Why Should We Care?", "The same tradeoff appears whenever human discussion creates reusable knowledge.", "Why care", ORANGE)
    xs = [0.7, 4.55, 8.4]
    titles = ["Course forums", "Developer forums", "Product communities"]
    bodies = [
        "Good student discussions help future batches. Too much AI can silently shrink that pool.",
        "Redirected traffic can overload expert communities like Stack Overflow-style forums.",
        "Automation can fix easy cases now while weakening the human knowledge pipeline later.",
    ]
    fills = [PALE, WHITE, PALE]
    for x, t, b, f in zip(xs, titles, bodies, fills):
        rrect(s, Inches(x), Inches(2.18), Inches(3.45), Inches(3.95), f, LINE, 0.08)
        txt(s, Inches(x + 0.22), Inches(2.45), Inches(2.6), Inches(0.22), t, size=19, bold=True, color=ORANGE if x < 1 else TEAL if x < 5 else GREEN)
        txt(s, Inches(x + 0.22), Inches(3.0), Inches(2.95), Inches(1.8), b, size=18, color=NAVY)
    rrect(s, Inches(1.55), Inches(6.32), Inches(10.2), Inches(0.40), WHITE, LINE, 0.08)
    txt(s, Inches(1.9), Inches(6.43), Inches(9.5), Inches(0.14), "Professor-friendly intuition: AI policy is also ecosystem policy.", size=18, bold=True, align=PP_ALIGN.CENTER)
    foot(s, "Refs: Horvitz (1999); Amershi et al. (2014); Bansal et al. (2021)", 3)


def s4_lit():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, TEAL)
    header(s, "Literature Review", "Three threads lead to our project.", "Literature", TEAL)
    data = [
        ("Selective response", TEAL, ["Chow (1970)", "El-Yaniv (2010)"], "When should the model abstain?"),
        ("Learning to defer", ORANGE, ["Madras et al. (2018)", "Mozannar and Sontag (2020)"], "When should humans take over?"),
        ("Congestion / commons", GREEN, ["Hardin (1968)", "Ostrom (1990)", "Kleinrock (1975)"], "What happens when shared systems get overloaded?"),
    ]
    for i, (title, accent, refs, q) in enumerate(data):
        x = 0.7 + 4*i
        rrect(s, Inches(x), Inches(2.18), Inches(3.55), Inches(3.95), WHITE, LINE, 0.08)
        txt(s, Inches(x + 0.22), Inches(2.45), Inches(2.8), Inches(0.24), title, size=19, bold=True, color=accent)
        txt(s, Inches(x + 0.22), Inches(3.0), Inches(2.7), Inches(0.8), "\n".join(refs), size=17)
        txt(s, Inches(x + 0.22), Inches(4.5), Inches(2.95), Inches(0.9), q, size=18, color=MUTED)
    rrect(s, Inches(1.2), Inches(6.25), Inches(10.9), Inches(0.45), PALE, LINE, 0.08)
    txt(s, Inches(1.45), Inches(6.37), Inches(10.3), Inches(0.15), "Gap: prior selective-response work does not treat the human forum as a congestible resource.", size=17, bold=True, align=PP_ALIGN.CENTER)
    foot(s, "Refs: Chow (1970); El-Yaniv (2010); Madras et al. (2018); Mozannar and Sontag (2020); Hardin (1968); Ostrom (1990)", 4)


def s5_base():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, ORANGE)
    header(s, "Base Paper", "Important: this idea is from the prior paper, not from us.", "Prior work", ORANGE)
    rrect(s, Inches(0.65), Inches(2.12), Inches(4.0), Inches(4.15), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.38), Inches(2.6), Inches(0.2), "What the base paper says", size=20, bold=True, color=ORANGE)
    bullets(s, Inches(0.88), Inches(2.78), Inches(3.3), Inches(2.6), [
        "AI should not always answer.",
        "Deferring some queries can improve long-run learning.",
        "Reason: human discussions create future signal."
    ], size=19)
    rrect(s, Inches(5.0), Inches(2.12), Inches(2.85), Inches(4.15), PALE, LINE, 0.08)
    txt(s, Inches(5.25), Inches(2.38), Inches(1.8), Inches(0.2), "Intuition", size=20, bold=True, color=ORANGE)
    rrect(s, Inches(5.55), Inches(3.05), Inches(1.75), Inches(0.7), WHITE, ORANGE, 0.18)
    rrect(s, Inches(5.55), Inches(4.55), Inches(1.75), Inches(0.7), WHITE, GREEN, 0.18)
    txt(s, Inches(5.96), Inches(3.28), Inches(0.95), Inches(0.14), "Always answer", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(5.86), Inches(4.78), Inches(1.15), Inches(0.14), "Selective response", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(5.2), Inches(5.58), Inches(2.0), Inches(0.4), "Sometimes less answering is better in the long run.", size=17, bold=True, align=PP_ALIGN.CENTER)
    rrect(s, Inches(8.15), Inches(2.12), Inches(4.7), Inches(4.15), PALE, LINE, 0.08)
    txt(s, Inches(8.42), Inches(2.38), Inches(2.8), Inches(0.2), "Why that is incomplete for us", size=20, bold=True, color=RED)
    bullets(s, Inches(8.38), Inches(2.78), Inches(4.0), Inches(2.6), [
        "Forum is not modeled as capacity-limited.",
        "No collapse dynamics under overload.",
        "No clean test of whether gains come from composition or just volume."
    ], size=18)
    foot(s, "Refs: Taitler et al. (2025)", 5)


def s6_our_method():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, TEAL)
    header(s, "Our Original Work", "We extend the base idea with congestion, capacity, and stronger validation checks.", "Our work", TEAL)
    rrect(s, Inches(0.65), Inches(2.12), Inches(6.0), Inches(4.15), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.38), Inches(2.0), Inches(0.2), "Method", size=20, bold=True, color=TEAL)
    bullets(s, Inches(0.88), Inches(2.75), Inches(3.0), Inches(2.4), [
        "Deterministic simulator, horizon T = 40.",
        "Tracks knowledge stock and forum capacity over time.",
        "Forum throughput falls under overload.",
        "Capacity can shrink if overload persists."
    ], size=18)
    rrect(s, Inches(4.1), Inches(3.0), Inches(1.1), Inches(0.68), PALE, TEAL, 0.16)
    rrect(s, Inches(5.55), Inches(3.0), Inches(1.1), Inches(0.68), PALE, ORANGE, 0.16)
    rrect(s, Inches(7.0), Inches(3.0), Inches(1.2), Inches(0.68), PALE, GREEN, 0.16)
    txt(s, Inches(4.38), Inches(3.22), Inches(0.55), Inches(0.14), "Users", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(5.83), Inches(3.22), Inches(0.55), Inches(0.14), "GenAI", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(7.32), Inches(3.22), Inches(0.55), Inches(0.14), "Forum", size=15, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(5.23), Inches(3.14), Inches(0.22), Inches(0.12), "->", size=21, bold=True, color=TEAL)
    txt(s, Inches(6.69), Inches(3.14), Inches(0.22), Inches(0.12), "->", size=21, bold=True, color=ORANGE)
    txt(s, Inches(4.15), Inches(4.55), Inches(3.6), Inches(0.35), "The human side is now a limited shared system.", size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    rrect(s, Inches(6.95), Inches(2.12), Inches(5.9), Inches(4.15), PALE, LINE, 0.08)
    txt(s, Inches(7.22), Inches(2.38), Inches(3.0), Inches(0.2), "What we additionally test", size=20, bold=True, color=TEAL)
    bullets(s, Inches(7.18), Inches(2.75), Inches(5.0), Inches(2.5), [
        "Matched answered-share control",
        "rho = 0 ablation",
        "Collapse boundary in (C1, xi) space",
        "Extended baselines: capacity-aware and novelty+capacity"
    ], size=18)
    foot(s, "Refs: project simulator outputs in sim/results and sim/figs", 6)


def s7_results():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, ORANGE)
    header(s, "Main Result", "Novelty-aware routing improves long-run learning across regimes.", "Results", ORANGE)
    rrect(s, Inches(0.65), Inches(2.12), Inches(3.4), Inches(4.15), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.38), Inches(2.2), Inches(0.2), "Say these numbers", size=20, bold=True, color=ORANGE)
    bullets(s, Inches(0.88), Inches(2.78), Inches(2.8), Inches(2.5), [
        "Stable: 26.15 -> 36.29",
        "Mid: 19.60 -> 31.56",
        "Collapse: 6.68 -> 9.54",
        "Mid overload also drops: 40 -> 25"
    ], size=18)
    txt(s, Inches(0.92), Inches(5.35), Inches(2.4), Inches(0.35), "Easy intuition: send better traffic to the forum.", size=16, color=MUTED)
    rrect(s, Inches(4.35), Inches(2.12), Inches(4.15), Inches(4.15), PALE, LINE, 0.08)
    pic(s, "sim/figs/extra/seed_robustness_ST.png", Inches(4.53), Inches(2.30), Inches(3.8), Inches(3.25))
    rrect(s, Inches(8.8), Inches(2.12), Inches(4.05), Inches(4.15), PALE, LINE, 0.08)
    pic(s, "sim/figs/extra/seed_robustness_overload.png", Inches(8.98), Inches(2.30), Inches(3.7), Inches(3.25))
    foot(s, "Refs: seed_robustness.csv", 7)


def s8_mechanism():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, RED)
    header(s, "Why We Believe The Result", "Two checks show the gain comes from composition, not just more answering.", "Validation", RED)
    rrect(s, Inches(0.65), Inches(2.12), Inches(3.1), Inches(4.15), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.38), Inches(2.0), Inches(0.2), "Check 1", size=20, bold=True, color=RED)
    txt(s, Inches(0.92), Inches(2.78), Inches(2.4), Inches(0.8), "Match the actual answered share.", size=20, bold=True)
    txt(s, Inches(0.92), Inches(3.72), Inches(2.35), Inches(1.2), "Even then, targeted still wins:\n31.35 vs 20.58", size=22, color=NAVY, bold=True)
    rrect(s, Inches(4.0), Inches(2.12), Inches(4.05), Inches(4.15), PALE, LINE, 0.08)
    pic(s, "sim/figs/extra/matched_answered_compare_ST.png", Inches(4.18), Inches(2.30), Inches(3.7), Inches(3.25))
    rrect(s, Inches(8.35), Inches(2.12), Inches(4.5), Inches(4.15), PALE, LINE, 0.08)
    txt(s, Inches(8.62), Inches(2.38), Inches(2.0), Inches(0.2), "Check 2", size=20, bold=True, color=RED)
    txt(s, Inches(8.62), Inches(2.78), Inches(3.5), Inches(0.75), "Set novelty value to zero: rho = 0.", size=20, bold=True)
    txt(s, Inches(8.62), Inches(3.70), Inches(3.6), Inches(0.65), "The advantage disappears.", size=24, color=NAVY, bold=True)
    pic(s, "sim/figs/extra/rho_ablation_gap_ST.png", Inches(9.05), Inches(4.15), Inches(3.35), Inches(1.40))
    foot(s, "Refs: matched_answered_compare.csv; rho_ablation.csv", 8)


def s9_baselines():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, GREEN)
    header(s, "Stronger Baselines", "Congestion feedback helps, but novelty+capacity is strongest in the mid regime.", "Baselines", GREEN)
    rrect(s, Inches(0.65), Inches(2.12), Inches(3.2), Inches(4.15), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.38), Inches(2.2), Inches(0.2), "Ranking", size=20, bold=True, color=GREEN)
    bullets(s, Inches(0.88), Inches(2.78), Inches(2.6), Inches(2.5), [
        "Uniform: 19.60",
        "Capacity-aware: 22.52",
        "Targeted novelty: 31.56",
        "Novelty + capacity: 33.23"
    ], size=18)
    txt(s, Inches(0.92), Inches(5.35), Inches(2.4), Inches(0.35), "Simple message: composition and congestion both matter.", size=15, color=MUTED)
    rrect(s, Inches(4.1), Inches(2.12), Inches(4.25), Inches(4.15), PALE, LINE, 0.08)
    pic(s, "sim/figs/extra/policy_compare_extended_ST.png", Inches(4.28), Inches(2.30), Inches(3.9), Inches(3.25))
    rrect(s, Inches(8.65), Inches(2.12), Inches(4.2), Inches(4.15), PALE, LINE, 0.08)
    pic(s, "sim/figs/extra/policy_compare_extended_overload.png", Inches(8.83), Inches(2.30), Inches(3.85), Inches(3.25))
    foot(s, "Refs: policy_compare_extended.csv", 9)


def s10_tipping():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, TEAL)
    header(s, "Tipping And Collapse", "The forum does not fail smoothly. There is a collapse region.", "Phase view", TEAL)
    rrect(s, Inches(0.65), Inches(2.12), Inches(3.15), Inches(4.15), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.38), Inches(2.0), Inches(0.2), "Intuition", size=20, bold=True, color=TEAL)
    bullets(s, Inches(0.88), Inches(2.78), Inches(2.55), Inches(2.6), [
        "If initial capacity is low, overload starts early.",
        "If overload damage is strong, recovery becomes hard.",
        "That creates a tipping boundary instead of a small gradual effect."
    ], size=18)
    rrect(s, Inches(4.0), Inches(2.12), Inches(4.2), Inches(4.15), PALE, LINE, 0.08)
    pic(s, "sim/figs/extra/collapse_boundary_c1_xi.png", Inches(4.18), Inches(2.30), Inches(3.85), Inches(3.25))
    rrect(s, Inches(8.55), Inches(2.12), Inches(4.3), Inches(4.15), PALE, LINE, 0.08)
    pic(s, "sim/figs/timeseries_collapse.png", Inches(8.73), Inches(2.30), Inches(3.95), Inches(3.25))
    foot(s, "Refs: collapse_boundary_c1_xi.csv", 10)


def s11_summary():
    s = prs.slides.add_slide(prs.slide_layouts[6]); base(s, ORANGE)
    header(s, "Summary And Future Work", "This is the slide to conclude the story clearly for class and professor.", "Summary", ORANGE)
    rrect(s, Inches(0.65), Inches(2.12), Inches(4.15), Inches(4.15), WHITE, LINE, 0.08)
    txt(s, Inches(0.92), Inches(2.38), Inches(2.1), Inches(0.2), "Three lines to conclude", size=20, bold=True, color=ORANGE)
    bullets(s, Inches(0.88), Inches(2.78), Inches(3.45), Inches(2.6), [
        "Base paper: selective response can help learning.",
        "Our extension: forums are congestible, so capacity matters.",
        "Main lesson: defer the right traffic, not just more traffic."
    ], size=18)
    rrect(s, Inches(5.1), Inches(2.12), Inches(3.5), Inches(4.15), PALE, LINE, 0.08)
    txt(s, Inches(5.37), Inches(2.38), Inches(1.8), Inches(0.2), "Future work", size=20, bold=True, color=ORANGE)
    bullets(s, Inches(5.33), Inches(2.78), Inches(2.8), Inches(2.4), [
        "Use real forum data for calibration.",
        "Learn policies instead of hand-designed rules.",
        "Model multiple forums and heterogeneous experts."
    ], size=18)
    rrect(s, Inches(8.9), Inches(2.12), Inches(3.95), Inches(4.15), PALE, LINE, 0.08)
    txt(s, Inches(9.17), Inches(2.38), Inches(1.8), Inches(0.2), "References", size=20, bold=True, color=TEAL)
    txt(s, Inches(9.15), Inches(2.78), Inches(3.15), Inches(2.8),
        "Taitler et al. (2025)\nChow (1970)\nEl-Yaniv (2010)\nMadras et al. (2018)\nMozannar and Sontag (2020)\nHardin (1968)\nOstrom (1990)\nVickrey (1969)", size=16)
    foot(s, "Refs: full bibliography in refs.bib", 11)


def thanks():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    rrect(s, Inches(0.9), Inches(1.0), Inches(5.2), Inches(4.8), WHITE, None, 0.08)
    txt(s, Inches(1.25), Inches(1.55), Inches(3.5), Inches(0.5), "Thank You", size=34, bold=True, color=NAVY)
    txt(s, Inches(1.27), Inches(2.35), Inches(2.0), Inches(0.25), "Questions?", size=24, bold=True, color=TEAL)
    txt(s, Inches(1.27), Inches(3.10), Inches(4.2), Inches(0.25), "Dikshant Gupta (24108)", size=18, color=NAVY)
    txt(s, Inches(1.27), Inches(3.45), Inches(4.2), Inches(0.25), "Paras Raina (24170)", size=18, color=NAVY)
    rrect(s, Inches(7.0), Inches(1.25), Inches(5.2), Inches(3.8), WHITE, None, 0.08)
    pic(s, "sim/figs/extra/policy_compare_extended_ST.png", Inches(7.25), Inches(1.48), Inches(4.7), Inches(3.0))
    txt(s, Inches(7.28), Inches(5.35), Inches(4.2), Inches(0.25), "Strongest mid-regime rule: novelty + capacity.", size=16, color=WHITE)


OUT.parent.mkdir(exist_ok=True)
title_slide()
s2_problem()
s3_examples()
s4_lit()
s5_base()
s6_our_method()
s7_results()
s8_mechanism()
s9_baselines()
s10_tipping()
s11_summary()
thanks()
prs.save(OUT)
print(OUT)
